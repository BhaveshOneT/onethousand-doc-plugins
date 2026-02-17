#!/usr/bin/env python3
"""
Hackathon Presentation Generator v2 — One Thousand GmbH

Generates branded hackathon presentations that EXACTLY match the original
manual format. All positions, font sizes, colors, and element types are
derived from the original PPTX analysis.

Key design decisions matching original:
- Theme colors (no explicit RGB) for most text — inherits from slide master
- Only #00B050 green highlights and specific #FFFFFF white text are explicit
- Manual TEXT_BOXes where original uses them (Pain, What's Next, etc.)
- Placeholder resizing where original resizes them
- AUTO_SHAPE with scheme fill for PoC summary lime box

Usage:
    python generate_v2.py \
        --template ot-hackathon-template.pptx \
        --variables variables.json \
        --content content.json \
        --output presentation.pptx \
        [--verbose]
"""

import argparse
import json
import logging
import re
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from copy import deepcopy

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from lxml import etree

# ---------------------------------------------------------------------------
# Brand Colors (from original PPTX XML analysis)
# ---------------------------------------------------------------------------
# Theme mapping: dk1=#000000, lt1=#FFFFFF, dk2=#242424 (ash), lt2=#D5F89E (lime)
# accent2=#19A960 (green)
# Most text uses THEME colors (inherited) — only these explicit colors needed:
OT_GREEN_HIGHLIGHT = RGBColor(0x00, 0xB0, 0x50)   # #00B050 — green inline highlights
OT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)               # Explicit white (only for What's Next content)
OT_BLACK = RGBColor(0x00, 0x00, 0x00)               # Black text
OT_MID_GRAY = RGBColor(0xBB, 0xBB, 0xBB)           # Placeholder hint text

# Note: OT_LIME_BG (#D5F89E) is applied via schemeClr "bg2" (=lt2), not explicit RGB
# Note: Most text does NOT set font.color.rgb — it inherits from theme

logging.basicConfig(level=logging.INFO, format='%(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

# ---------------------------------------------------------------------------
# Rich Text Parsing — **bold** and <<green>>
# ---------------------------------------------------------------------------
_RICH_RE = re.compile(r'(\*\*.*?\*\*|<<.*?>>)')


def _parse_rich_segments(text: str) -> List[Tuple[str, str]]:
    """Parse text with **bold** and <<green>> markup into segments."""
    parts = _RICH_RE.split(text)
    segments = []
    for part in parts:
        if not part:
            continue
        if part.startswith('**') and part.endswith('**'):
            segments.append((part[2:-2], 'bold'))
        elif part.startswith('<<') and part.endswith('>>'):
            segments.append((part[2:-2], 'green'))
        else:
            segments.append((part, 'normal'))
    return segments


def add_rich_paragraph(text_frame, text: str, font_size: int = 13,
                       use_theme_color: bool = True,
                       explicit_color: RGBColor = None,
                       green_color: RGBColor = OT_GREEN_HIGHLIGHT,
                       spacing_before: int = 6, spacing_after: int = 6,
                       is_first: bool = False, alignment=None,
                       bold_base: bool = False):
    """Add a paragraph with mixed bold/green/normal formatting.

    Markup:
      **text** → bold, theme or explicit color
      <<text>> → normal, green_color (#00B050)
      plain    → normal, theme or explicit color

    If use_theme_color=True, normal/bold runs do NOT set font.color.rgb
    (they inherit from theme). Only <<green>> runs get explicit color.
    If explicit_color is set, it overrides theme for normal/bold runs.
    """
    p = text_frame.paragraphs[0] if is_first else text_frame.add_paragraph()
    p.space_before = Pt(spacing_before)
    p.space_after = Pt(spacing_after)
    if alignment:
        p.alignment = alignment

    segments = _parse_rich_segments(text)
    for seg_text, style in segments:
        r = p.add_run()
        r.text = seg_text
        r.font.size = Pt(font_size)
        if style == 'bold':
            r.font.bold = True
            if explicit_color:
                r.font.color.rgb = explicit_color
            # else: theme color (no explicit set)
        elif style == 'green':
            r.font.bold = False
            r.font.color.rgb = green_color
        else:
            r.font.bold = bold_base
            if explicit_color:
                r.font.color.rgb = explicit_color
            # else: theme color (no explicit set)
    return p


# ---------------------------------------------------------------------------
# Utility helpers
# ---------------------------------------------------------------------------

def load_json(path: Path) -> Dict[str, Any]:
    with open(path, 'r', encoding='utf-8') as f:
        return json.load(f)


def find_layout(prs, name: str):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    logger.warning(f"Layout not found: {name}")
    return None


def ph_by_idx(slide, idx: int):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def set_ph_text_theme(ph, text: str, size: int = None, bold: bool = False):
    """Set text on a placeholder using THEME color (no explicit RGB).
    This matches the original PPTX behavior where most text inherits from theme.
    """
    if ph is None or not ph.has_text_frame:
        return
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    if size:
        r.font.size = Pt(size)
    if bold:
        r.font.bold = True
    # NO font.color.rgb set — inherits theme color


def set_ph_text(ph, text: str, size: int = None, bold: bool = False,
                color: RGBColor = None):
    """Set text on a placeholder with optional explicit color."""
    if ph is None or not ph.has_text_frame:
        return
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    if size:
        r.font.size = Pt(size)
    if bold:
        r.font.bold = True
    if color:
        r.font.color.rgb = color


def reposition_shape(shape, left=None, top=None, width=None, height=None):
    """Reposition/resize a shape using inches."""
    if left is not None:
        shape.left = Inches(left)
    if top is not None:
        shape.top = Inches(top)
    if width is not None:
        shape.width = Inches(width)
    if height is not None:
        shape.height = Inches(height)


def add_textbox_theme(slide, left, top, width, height, text,
                      font_size=14, alignment=PP_ALIGN.LEFT, bold=False,
                      word_wrap=True, vertical_anchor=MSO_ANCHOR.TOP):
    """Add a textbox using THEME color (no explicit RGB set on font).
    Text inherits color from the slide master/theme.
    """
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = vertical_anchor
    p = tf.paragraphs[0]
    p.alignment = alignment
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    # NO font.color.rgb — inherits from theme
    return tb


def add_textbox(slide, left, top, width, height, text,
                font_size=14, text_color=None, alignment=PP_ALIGN.LEFT,
                bold=False, word_wrap=True, vertical_anchor=MSO_ANCHOR.TOP,
                bg_color=None):
    """Add a textbox with explicit color."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = vertical_anchor
    p = tf.paragraphs[0]
    p.alignment = alignment
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    if text_color:
        r.font.color.rgb = text_color
    if bg_color:
        tb.fill.solid()
        tb.fill.fore_color.rgb = bg_color
    return tb


def add_rich_textbox_theme(slide, left, top, width, height, text: str,
                           font_size=14, green_color=OT_GREEN_HIGHLIGHT,
                           alignment=PP_ALIGN.LEFT, word_wrap=True,
                           vertical_anchor=MSO_ANCHOR.TOP,
                           bold_base=False):
    """Add a textbox with rich text markup, using THEME colors for base text."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = vertical_anchor

    add_rich_paragraph(tf, text, font_size=font_size,
                       use_theme_color=True,
                       green_color=green_color,
                       is_first=True, alignment=alignment,
                       bold_base=bold_base)
    return tb


def add_auto_shape_with_scheme_fill(slide, left, top, width, height,
                                     scheme_color="bg2"):
    """Add an AUTO_SHAPE (RECTANGLE) with scheme color fill.
    Original PoC summary uses schemeClr "bg2" which maps to lt2=#D5F89E (lime).
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    # Remove default outline
    shape.line.fill.background()

    # Set scheme color fill via XML manipulation
    spPr = shape._element.find(f'{{{NS_A}}}spPr')
    if spPr is not None:
        # Remove existing fill
        for child in list(spPr):
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag in ('solidFill', 'noFill', 'gradFill', 'pattFill'):
                spPr.remove(child)
        # Add scheme color fill
        solidFill = etree.SubElement(spPr, f'{{{NS_A}}}solidFill')
        schemeClr = etree.SubElement(solidFill, f'{{{NS_A}}}schemeClr')
        schemeClr.set('val', scheme_color)

    return shape


def image_placeholder(slide, left, top, width, height, description):
    """Add a subtle image placeholder — no heavy bg fill, just hint text."""
    add_textbox(slide, left, top, width, height,
                f"[IMAGE: {description}]",
                font_size=14, text_color=OT_MID_GRAY,
                alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE)


def add_speaker_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def set_footer_textbox(slide, slide_num, copyright_text="© 2019-2026 ONE THOUSAND"):
    """Add footer as textbox matching original format: center-aligned at bottom."""
    add_textbox(slide, 4.62, 7.05, 4.09, 0.20, copyright_text,
                font_size=8, text_color=OT_MID_GRAY,
                alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.TOP)


def slide_count(prs):
    return len(prs.slides)


def _enable_autofit(text_frame):
    """Enable auto-shrink text fitting — <a:normAutofit/> without fontScale."""
    body_props = text_frame._txBody.find(f'{{{NS_A}}}bodyPr')
    if body_props is None:
        return
    for child in list(body_props):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('noAutofit', 'normAutofit', 'spAutoFit'):
            body_props.remove(child)
    etree.SubElement(body_props, f'{{{NS_A}}}normAutofit')


def remove_shape_by_ph_idx(slide, idx):
    """Remove a placeholder element entirely from the slide XML."""
    ph = ph_by_idx(slide, idx)
    if ph is not None:
        sp = ph._element
        sp.getparent().remove(sp)


# ---------------------------------------------------------------------------
# Slide creation functions — EXACT match to original PPTX
# ---------------------------------------------------------------------------

def make_cover(prs, client_name, location, date, use_case_title, day=1, verbose=False):
    """Cover slide — 'Title Lime + one Logo' layout.

    Original formatting (from uploaded PPTX analysis):
    - PH idx=10 REPOSITIONED to (0.29, 6.05) 3.90x0.23 — date/location
    - PH idx=0 REPOSITIONED to (0.29, 6.41) 9.33x1.46 — TWO paragraphs
    - PH idx=11 (picture) and PH idx=15 (text) REMOVED — not used by original
    - PICTURE at (9.60, 6.55) — Client logo (manual, not from layout)
    """
    layout = find_layout(prs, "Title Lime + one Logo")
    if not layout:
        layout = find_layout(prs, "Cover Lime + one Logo")
        if not layout:
            return
    slide = prs.slides.add_slide(layout)

    # REMOVE unused layout placeholders that cause "duplicate" appearance
    remove_shape_by_ph_idx(slide, 11)   # Picture placeholder from layout
    remove_shape_by_ph_idx(slide, 15)   # Text placeholder from layout

    # Date/location — PH idx=10, reposition to match original
    date_text = f"{location}  |  {date}"
    if day == 2:
        date_text += "  |  Day 2"
    date_ph = ph_by_idx(slide, 10)
    if date_ph:
        reposition_shape(date_ph, left=0.29, top=6.05, width=3.90, height=0.23)
        set_ph_text_theme(date_ph, date_text)
    else:
        # Fallback: create textbox if PH 10 not in layout
        add_textbox_theme(slide, 0.29, 6.05, 3.90, 0.23, date_text,
                          font_size=10, alignment=PP_ALIGN.LEFT)

    # Title — PH idx=0, reposition to match original (0.29, 6.41) 9.33x1.46
    title_ph = ph_by_idx(slide, 0)
    if title_ph:
        reposition_shape(title_ph, left=0.29, top=6.41, width=9.33, height=1.46)
        if title_ph.has_text_frame:
            tf = title_ph.text_frame
            tf.clear()
            # Paragraph 1: Title — 32pt bold, theme color
            p1 = tf.paragraphs[0]
            r1 = p1.add_run()
            r1.text = f"Strengthening {client_name} With AI"
            r1.font.size = Pt(32)
            r1.font.bold = True

            # Paragraph 2: Subtitle — 32pt NOT bold, theme color
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = f"AI Hackathon | {use_case_title}"
            r2.font.size = Pt(32)
            r2.font.bold = False

    # Client logo placeholder at bottom-right
    image_placeholder(slide, 9.60, 6.55, 4.18, 0.84, "Client logo")

    if verbose:
        logger.info(f"Created Day {day} cover slide")


def make_checkin(prs, questions, client_name, verbose=False):
    """Check-in slide — 'Bullet Points Ash' layout.

    Original formatting:
    - PH idx=0 REPOSITIONED to (0.44, 6.26) 10.02x0.91 — "Check-in" at 66pt, BOTTOM
    - TEXT_BOX at (4.53, 1.28) 8.62x0.81 — Questions at 32pt, theme color
    - ALL emoji placeholders REMOVED entirely
    """
    layout = find_layout(prs, "Bullet Points Ash")
    if not layout:
        return
    slide = prs.slides.add_slide(layout)

    questions = [q.replace("{client_name}", client_name) for q in questions]

    # Remove ALL emoji placeholders (idx=27, 46, 47)
    for emoji_idx in [27, 46, 47]:
        remove_shape_by_ph_idx(slide, emoji_idx)

    # Clear unused column placeholders
    for idx in [1, 36, 42, 44, 45]:
        ph = ph_by_idx(slide, idx)
        if ph:
            remove_shape_by_ph_idx(slide, idx)

    # PH idx=0 — Reposition to BOTTOM of slide, 66pt, Wavetable font
    title_ph = ph_by_idx(slide, 0)
    if title_ph:
        reposition_shape(title_ph, left=0.44, top=6.26, width=10.02, height=0.91)
        if title_ph.has_text_frame:
            tf = title_ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = "Check-in"
            r.font.size = Pt(66)
            r.font.name = "Wavetable"

    # Questions — single TEXT_BOX at (4.53, 1.28) 8.62x0.81
    # EXPLICIT white needed: textboxes don't inherit theme on dark Ash bg
    if questions:
        tb = slide.shapes.add_textbox(
            Inches(4.53), Inches(1.28), Inches(8.62), Inches(0.81)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        _enable_autofit(tf)
        for i, q in enumerate(questions):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = q
            r.font.size = Pt(32)
            r.font.color.rgb = OT_WHITE  # Explicit white on dark bg

    set_footer_textbox(slide, slide_count(prs))
    add_speaker_notes(slide, "Engage the team with check-in questions.")
    if verbose:
        logger.info("Created Check-in slide")


def make_agenda(prs, agenda_data, verbose=False):
    """Agenda slide — 'Dayline Lime' layout."""
    layout = find_layout(prs, "Dayline Lime")
    if not layout:
        return
    slide = prs.slides.add_slide(layout)

    set_ph_text_theme(ph_by_idx(slide, 0), "AI Hackathon – Agenda", size=36, bold=True)
    set_ph_text(ph_by_idx(slide, 1), "")
    set_ph_text_theme(ph_by_idx(slide, 21), "Day 1")
    set_ph_text_theme(ph_by_idx(slide, 121), "Day 2")

    day1 = agenda_data.get("day1", [])
    day2 = agenda_data.get("day2", [])

    # Day 1 grid
    d1_grid = [(97, 107, 108), (66, 115, 116), (66, 117, 118)]
    for row_idx, (t_idx, a1_idx, a2_idx) in enumerate(d1_grid):
        if row_idx < len(day1):
            set_ph_text_theme(ph_by_idx(slide, t_idx), day1[row_idx].get("time", ""))
            set_ph_text_theme(ph_by_idx(slide, a1_idx), day1[row_idx].get("activity", ""))
            set_ph_text(ph_by_idx(slide, a2_idx), "")
        else:
            set_ph_text(ph_by_idx(slide, t_idx), "")
            set_ph_text(ph_by_idx(slide, a1_idx), "")
            set_ph_text(ph_by_idx(slide, a2_idx), "")

    for idx in [109, 110, 111, 112, 113, 114, 119, 120]:
        set_ph_text(ph_by_idx(slide, idx), "")

    d2_grid = [(122, 124, 125), (123, 132, 138)]
    for row_idx, (t_idx, a1_idx, a2_idx) in enumerate(d2_grid):
        if row_idx < len(day2):
            set_ph_text_theme(ph_by_idx(slide, t_idx), day2[row_idx].get("time", ""))
            set_ph_text_theme(ph_by_idx(slide, a1_idx), day2[row_idx].get("activity", ""))
            set_ph_text(ph_by_idx(slide, a2_idx), "")
        else:
            set_ph_text(ph_by_idx(slide, t_idx), "")
            set_ph_text(ph_by_idx(slide, a1_idx), "")
            set_ph_text(ph_by_idx(slide, a2_idx), "")

    for idx in [126, 129, 130, 131, 134, 135, 136]:
        set_ph_text(ph_by_idx(slide, idx), "")

    set_ph_text(ph_by_idx(slide, 22), "")
    if verbose:
        logger.info("Created Agenda slide")


def make_toc(prs, verbose=False):
    """Table of Contents — 'Table of Contents large' layout.

    Original keeps only: number PHs (37,40,20,32) and title PHs (36,39,19,31).
    Arrow PHs (38,41,21,33) and description PHs (34,35,42,43) are REMOVED.
    """
    layout = find_layout(prs, "Table of Contents large")
    if not layout:
        return
    slide = prs.slides.add_slide(layout)

    # Clear title PH
    set_ph_text(ph_by_idx(slide, 0), "")

    # Fill TOC items — (number_idx, title_idx)
    items = [("01", "Pain", 37, 36), ("02", "Data", 40, 39),
             ("03", "APPROACH", 20, 19), ("04", "Challenges", 32, 31)]

    for num, title, n_idx, h_idx in items:
        set_ph_text_theme(ph_by_idx(slide, n_idx), num)
        set_ph_text_theme(ph_by_idx(slide, h_idx), title)

    # REMOVE arrow placeholders entirely (not in original)
    for idx in [38, 41, 21, 33]:
        remove_shape_by_ph_idx(slide, idx)

    # REMOVE description area placeholders entirely (empty in original)
    for idx in [34, 35, 42, 43]:
        remove_shape_by_ph_idx(slide, idx)

    set_footer_textbox(slide, slide_count(prs))
    if verbose:
        logger.info("Created TOC slide")


def make_pain(prs, bullets, verbose=False):
    """Pain slide — 'Chapter Divider Ash + Text' layout.

    Original formatting:
    - PH idx=14 — "01" inherited theme
    - PH idx=0 — "Pain" inherited theme (default position)
    - TEXT_BOX (NOT PH!) at (4.26, 3.23) 8.70x3.49 — 18pt with #00B050 green highlights
    - Footer TEXT_BOX at (4.62, 7.05) center-aligned
    """
    layout = find_layout(prs, "Chapter Divider Ash + Text")
    if not layout:
        return
    slide = prs.slides.add_slide(layout)

    # PH idx=14 — chapter number "01", theme color
    set_ph_text_theme(ph_by_idx(slide, 14), "01")

    # PH idx=0 — chapter title "Pain", theme color
    set_ph_text_theme(ph_by_idx(slide, 0), "Pain")

    # Remove the default PH idx=15 — we use a manual TEXT_BOX instead
    remove_shape_by_ph_idx(slide, 15)

    # Content — manual TEXT_BOX at exact original position
    # EXPLICIT white needed: textboxes don't inherit theme on dark Ash bg
    tb = slide.shapes.add_textbox(
        Inches(4.26), Inches(3.23), Inches(8.70), Inches(3.49)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, bullet in enumerate(bullets):
        add_rich_paragraph(tf, bullet, font_size=18,
                           use_theme_color=False,
                           explicit_color=OT_WHITE,
                           green_color=OT_GREEN_HIGHLIGHT,
                           spacing_before=4, spacing_after=4,
                           is_first=(i == 0),
                           alignment=PP_ALIGN.LEFT)

    set_footer_textbox(slide, slide_count(prs))
    if verbose:
        logger.info("Created Pain slide")


def make_data(prs, bullets, verbose=False):
    """Data slide — 'Chapter Divider Ash + Text' layout.

    Original formatting:
    - PH idx=14 — "02" theme
    - PH idx=0 — "Data" theme
    - PH idx=15 RESIZED to (4.15, 3.39) 8.70x2.76 — 18pt theme, bold prefixes
    """
    layout = find_layout(prs, "Chapter Divider Ash + Text")
    if not layout:
        return
    slide = prs.slides.add_slide(layout)

    set_ph_text_theme(ph_by_idx(slide, 14), "02")
    set_ph_text_theme(ph_by_idx(slide, 0), "Data")

    # Resize PH idx=15 to match original
    content_ph = ph_by_idx(slide, 15)
    if content_ph:
        reposition_shape(content_ph, left=4.15, top=3.39, width=8.70, height=2.76)
        if content_ph.has_text_frame:
            tf = content_ph.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.TOP
            _enable_autofit(tf)

            # Original uses: " " + "Bold title: " accent2 + "description" theme
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_before = Pt(4)
                p.space_after = Pt(4)
                p.alignment = PP_ALIGN.LEFT
                # Parse **bold** markup
                segments = _parse_rich_segments(bullet)
                # Leading space like original
                r_sp = p.add_run()
                r_sp.text = " "
                r_sp.font.size = Pt(18)
                for seg_text, style in segments:
                    r = p.add_run()
                    r.text = seg_text
                    r.font.size = Pt(18)
                    if style == 'bold':
                        r.font.bold = True
                        _set_run_scheme_color(r, "accent2")  # Green accent
                    elif style == 'green':
                        r.font.color.rgb = OT_GREEN_HIGHLIGHT
                    # else: theme color (inherited white on dark bg)

    set_footer_textbox(slide, slide_count(prs))
    if verbose:
        logger.info("Created Data slide")


def make_data_screenshots(prs, verbose=False):
    """Data screenshots slide — 'Title Ash + small Image' layout."""
    layout = find_layout(prs, "Title Ash + small Image")
    if not layout:
        return
    slide = prs.slides.add_slide(layout)

    set_ph_text_theme(ph_by_idx(slide, 0), "Data Screenshots", size=24)
    set_ph_text_theme(ph_by_idx(slide, 32), "Sample data files: emails, PDFs, spreadsheets", size=12)
    set_ph_text(ph_by_idx(slide, 31), "")

    set_footer_textbox(slide, slide_count(prs))
    if verbose:
        logger.info("Created Data Screenshots slide")


def make_approach(prs, bullets, verbose=False):
    """Approach slide — 'Chapter Divider Ash + Text' layout.

    Original formatting:
    - PH idx=14 — "03" theme
    - PH idx=0 RESIZED to (0.29, 3.27) 3.67x0.66 — "APPROACH" at 48pt theme
    - PH idx=15 RESIZED to (4.35, 3.11) 8.90x3.88 — 16pt, anchor=t, LEFT, bold prefixes
    """
    layout = find_layout(prs, "Chapter Divider Ash + Text")
    if not layout:
        return
    slide = prs.slides.add_slide(layout)

    set_ph_text_theme(ph_by_idx(slide, 14), "03")

    # Title — resize and set 48pt
    title_ph = ph_by_idx(slide, 0)
    if title_ph:
        reposition_shape(title_ph, left=0.29, top=3.27, width=3.67, height=0.66)
        set_ph_text_theme(title_ph, "APPROACH", size=48)

    # Content — resize PH idx=15
    content_ph = ph_by_idx(slide, 15)
    if content_ph:
        reposition_shape(content_ph, left=4.35, top=3.11, width=8.90, height=3.88)
        if content_ph.has_text_frame:
            tf = content_ph.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.TOP
            _enable_autofit(tf)

            # Original: bold titles in ACCENT_2, descriptions in theme,
            # last item (question) uses BACKGROUND_2 for bold
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_before = Pt(7.5)
                p.alignment = PP_ALIGN.LEFT
                is_question = (i == len(bullets) - 1)  # Last item is the question
                segments = _parse_rich_segments(bullet)
                for seg_text, style in segments:
                    r = p.add_run()
                    r.text = seg_text
                    r.font.size = Pt(16)
                    if style == 'bold':
                        r.font.bold = True
                        if is_question:
                            _set_run_scheme_color(r, "bg2")  # Lime for question
                        else:
                            _set_run_scheme_color(r, "accent2")  # Green accent
                    elif style == 'green':
                        r.font.color.rgb = OT_GREEN_HIGHLIGHT
                    else:
                        if is_question:
                            # Question description also uses bg2
                            pass  # theme inherited
                        # else: theme color (inherited)

    set_footer_textbox(slide, slide_count(prs))
    if verbose:
        logger.info("Created Approach slide")


def make_challenges(prs, bullets, verbose=False):
    """Challenges slide — 'Chapter Divider Ash + Text' layout.

    Original formatting:
    - PH idx=14 — "04" theme
    - PH idx=0 RESIZED to (0.29, 3.27) 4.46x0.66 — "Challenges" at 48pt theme
    - PH idx=15 RESIZED to (4.76, 3.27) 7.91x3.65 — 14pt, anchor=t, bold prefixes
    """
    layout = find_layout(prs, "Chapter Divider Ash + Text")
    if not layout:
        return
    slide = prs.slides.add_slide(layout)

    set_ph_text_theme(ph_by_idx(slide, 14), "04")

    title_ph = ph_by_idx(slide, 0)
    if title_ph:
        reposition_shape(title_ph, left=0.29, top=3.27, width=4.46, height=0.66)
        set_ph_text_theme(title_ph, "Challenges", size=48)

    content_ph = ph_by_idx(slide, 15)
    if content_ph:
        reposition_shape(content_ph, left=4.76, top=3.27, width=7.91, height=3.65)
        if content_ph.has_text_frame:
            tf = content_ph.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.TOP
            _enable_autofit(tf)

            # Original: bold titles in ACCENT_2, " " spacer, description in theme
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_before = Pt(7.5)
                p.alignment = PP_ALIGN.LEFT
                segments = _parse_rich_segments(bullet)
                for seg_text, style in segments:
                    r = p.add_run()
                    r.text = seg_text
                    r.font.size = Pt(14)
                    if style == 'bold':
                        r.font.bold = True
                        _set_run_scheme_color(r, "accent2")  # Green accent
                    elif style == 'green':
                        r.font.color.rgb = OT_GREEN_HIGHLIGHT
                    # else: theme color (inherited)

    set_footer_textbox(slide, slide_count(prs))
    if verbose:
        logger.info("Created Challenges slide")


def make_divider(prs, title, verbose=False):
    """Green chapter divider — 'Chapter Divider Lime' layout."""
    layout = find_layout(prs, "Chapter Divider Lime")
    if not layout:
        return
    slide = prs.slides.add_slide(layout)

    set_ph_text_theme(ph_by_idx(slide, 0), title, size=44, bold=True)
    set_ph_text(ph_by_idx(slide, 14), "")
    set_ph_text(ph_by_idx(slide, 32), "")

    if verbose:
        logger.info(f"Created divider: {title}")


def make_process_flow(prs, verbose=False):
    """Process Flow slide — DEFAULT layout.

    Original formatting:
    - TEXT_BOX at (0.23, 0.65) 12.82x1.35 — "WE DISCUSSED THE PROCESS FLOW" 40pt THEME color
    - PICTURES below for user to add
    """
    layout = find_layout(prs, "DEFAULT")
    if not layout:
        return
    slide = prs.slides.add_slide(layout)

    # Title — 40pt, EXPLICIT white (textboxes on dark DEFAULT bg need explicit color)
    add_textbox(slide, 0.23, 0.65, 12.82, 1.35,
                "WE DISCUSSED THE PROCESS FLOW",
                font_size=40, text_color=OT_WHITE, bold=False,
                alignment=PP_ALIGN.LEFT,
                vertical_anchor=MSO_ANCHOR.TOP)

    # Image placeholder
    image_placeholder(slide, 0.5, 2.5, 12.0, 4.5,
                      "Add process flow diagram (from whiteboard / Miro)")

    set_footer_textbox(slide, slide_count(prs))
    if verbose:
        logger.info("Created Process Flow slide")


def make_architecture(prs, verbose=False):
    """Architecture slide — DEFAULT layout.

    Original formatting:
    - TEXT_BOX at (0.46, 0.85) 6.21x2.42 — "WE'VE SET UP AN INITIAL ARCHITECTURE" 50pt THEME color
    - Image area for architecture diagram
    """
    layout = find_layout(prs, "DEFAULT")
    if not layout:
        return
    slide = prs.slides.add_slide(layout)

    # Title — 50pt, EXPLICIT white (textboxes on dark DEFAULT bg)
    add_textbox(slide, 0.46, 0.85, 6.21, 2.42,
                "WE'VE SET UP AN INITIAL ARCHITECTURE",
                font_size=50, text_color=OT_WHITE, bold=False,
                alignment=PP_ALIGN.LEFT,
                vertical_anchor=MSO_ANCHOR.TOP)

    # Image placeholder
    image_placeholder(slide, 0.5, 3.5, 12.0, 3.5,
                      "Add architecture diagram (from Miro / whiteboard)")

    set_footer_textbox(slide, slide_count(prs))
    if verbose:
        logger.info("Created Architecture slide")


def make_business_value(prs, bv_data, verbose=False):
    """Business value slide — 'Table of Contents small' layout.

    Original formatting:
    - TEXT_BOX (NOT PH!) at (0.23, 1.21) 12.82x1.35 — Title 20pt theme
    - PH idx=13 anchor=ctr — item 1: 16pt bold title + 16pt description, theme
    - PH idx=15 anchor=ctr — item 2
    - PH idx=19 anchor=ctr — item 3
    - PH idx=14, 16, 20 — Numbers "01", "02", "03"
    """
    layout = find_layout(prs, "Table of Contents small")
    if not layout:
        return
    slide = prs.slides.add_slide(layout)

    # Clear default title PH idx=0 — we use a manual textbox instead
    title_ph = ph_by_idx(slide, 0)
    if title_ph:
        set_ph_text(title_ph, "")

    # Title — manual TEXT_BOX with "CREATE BUSINESS VALUE" highlighted
    title_text = bv_data.get("title",
        "THE OVERALL GOAL IS TO <<CREATE BUSINESS VALUE>>")
    # Build textbox with scheme color highlight
    tb = slide.shapes.add_textbox(
        Inches(0.23), Inches(1.21), Inches(12.82), Inches(1.35)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    segments = _parse_rich_segments(title_text)
    for seg_text, style in segments:
        r = p.add_run()
        r.text = seg_text
        r.font.size = Pt(20)
        if style == 'green':
            # <<green>> in title = scheme color highlight (tx1 = dark on this layout)
            _set_run_scheme_color(r, "tx1")
        else:
            r.font.color.rgb = OT_WHITE  # Normal text white on dark bg

    items = bv_data.get("items", [])
    # Number placeholders and content placeholders
    slot_mapping = [(14, 13), (16, 15), (20, 19)]

    for i, item in enumerate(items[:3]):
        if i < len(slot_mapping):
            num_idx, content_idx = slot_mapping[i]
            num = item.get("number", f"{i+1:02d}")
            item_title = item.get("title", "")
            desc = item.get("description", "")

            # Number — theme color
            set_ph_text_theme(ph_by_idx(slide, num_idx), num)

            # Content — single paragraph: "Bold Title" + " " + "- description" + "green end"
            # Matching original format exactly
            content_ph = ph_by_idx(slide, content_idx)
            if content_ph and content_ph.has_text_frame:
                tf = content_ph.text_frame
                tf.clear()
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                _enable_autofit(tf)

                p1 = tf.paragraphs[0]
                # Bold title
                r1 = p1.add_run()
                r1.text = item_title
                r1.font.bold = True
                r1.font.size = Pt(16)

                if desc:
                    # Parse description for <<green>> highlights
                    plain_desc = re.sub(r'<<|>>', '', desc)
                    segments = _parse_rich_segments(desc)
                    # Space separator
                    r_sp = p1.add_run()
                    r_sp.text = " "
                    r_sp.font.size = Pt(16)
                    # "- description" with green highlights
                    for seg_text, style in segments:
                        r = p1.add_run()
                        r.font.size = Pt(16)
                        if style == 'green':
                            r.text = seg_text
                            _set_run_scheme_color(r, "bg2")  # Lime green
                        elif style == 'bold':
                            r.text = seg_text
                            r.font.bold = True
                        else:
                            r.text = seg_text

    # Clear unused arrow/extra placeholders
    for idx in [17, 18, 21]:
        set_ph_text(ph_by_idx(slide, idx), "")

    set_footer_textbox(slide, slide_count(prs))
    if verbose:
        logger.info("Created Business Value slide")


def _set_run_scheme_color(run, scheme_val="tx1"):
    """Set a run's font color using scheme color (e.g., "tx1" for dark text).
    This matches the original PPTX which uses schemeClr "tx1" for text on lime bg.
    """
    rPr = run._r.get_or_add_rPr()
    # Remove any existing color
    for child in list(rPr):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag == 'solidFill':
            rPr.remove(child)
    solidFill = etree.SubElement(rPr, f'{{{NS_A}}}solidFill')
    schemeClr = etree.SubElement(solidFill, f'{{{NS_A}}}schemeClr')
    schemeClr.set('val', scheme_val)


def make_poc_summary(prs, intro, features, verbose=False):
    """PoC Summary slide — DEFAULT layout.

    Original formatting:
    - TEXT_BOX at (0.23, 0.65) 12.82x1.35 — title 40pt, white on dark bg
    - AUTO_SHAPE (RECTANGLE) at (3.55, 1.96) 6.24x4.37 — scheme fill "bg2" (lime)
      Text uses schemeClr "tx1" (dark/black) for readability on lime bg
    """
    layout = find_layout(prs, "DEFAULT")
    if not layout:
        return
    slide = prs.slides.add_slide(layout)

    # Title — 40pt, EXPLICIT white (textbox on dark DEFAULT bg)
    add_textbox(slide, 0.23, 0.65, 12.82, 1.35,
                "WE HAVE FOCUSED ON THE CORE PAIN POINTS",
                font_size=40, text_color=OT_WHITE, bold=False,
                alignment=PP_ALIGN.LEFT,
                vertical_anchor=MSO_ANCHOR.TOP)

    # Lime rectangle with scheme fill "bg2" (=lt2=#D5F89E)
    rect = add_auto_shape_with_scheme_fill(
        slide, 3.55, 1.96, 6.24, 4.37, scheme_color="bg2"
    )

    # Add text to the rectangle shape
    tf = rect.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # First line: intro text — with line_spacing=1.5 like original
    p0 = tf.paragraphs[0]
    p0.line_spacing = 1.5
    p0.alignment = PP_ALIGN.LEFT
    r0 = p0.add_run()
    r0.text = intro
    r0.font.size = Pt(16)
    r0.font.bold = False
    _set_run_scheme_color(r0, "tx1")  # Dark text on lime bg

    # Feature lines — bold, dark on lime, with line_spacing=1.5
    for feature in features:
        p = tf.add_paragraph()
        p.line_spacing = 1.5
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = feature
        r.font.size = Pt(16)
        r.font.bold = True
        _set_run_scheme_color(r, "tx1")  # Dark text on lime bg

    set_footer_textbox(slide, slide_count(prs))
    if verbose:
        logger.info("Created PoC Summary slide")


def make_image_slide_blank(prs, image_desc, verbose=False):
    """Blank image slide for team photos — DEFAULT layout."""
    layout = find_layout(prs, "DEFAULT")
    if not layout:
        return
    slide = prs.slides.add_slide(layout)

    image_placeholder(slide, 0.5, 0.5, 12.0, 6.5, image_desc)
    set_footer_textbox(slide, slide_count(prs))
    if verbose:
        logger.info(f"Created blank image slide: {image_desc}")


def make_whats_next(prs, next_steps, verbose=False):
    """What's Next slide — 'Chapter Divider Ash + Text' layout.

    Original formatting:
    - PH idx=0, 14, 15 are REMOVED entirely (not used)
    - TEXT_BOX at (0.23, 0.65) 12.82x1.35 — "WHAT'S NEXT?" 40pt white
    - TEXT_BOX at (0.42, 2.18) 12.62x3.32 — Content at 28pt:
      #FFFFFF for normal text, bold+white for key phrases
    """
    layout = find_layout(prs, "Chapter Divider Ash + Text")
    if not layout:
        return
    slide = prs.slides.add_slide(layout)

    # REMOVE all default placeholders entirely — original doesn't use them
    for idx in [0, 14, 15]:
        remove_shape_by_ph_idx(slide, idx)

    # Title — manual TEXT_BOX, 40pt, EXPLICIT white (dark Ash bg)
    add_textbox(slide, 0.23, 0.65, 12.82, 1.35,
                "WHAT'S NEXT?", font_size=40, text_color=OT_WHITE,
                bold=False, alignment=PP_ALIGN.LEFT,
                vertical_anchor=MSO_ANCHOR.TOP)

    # Content — manual TEXT_BOX at (0.42, 2.18) 12.62x3.32
    tb = slide.shapes.add_textbox(
        Inches(0.42), Inches(2.18), Inches(12.62), Inches(3.32)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    # Original starts content at P1 (P0 empty), uses sp_after=6, line_sp=1.0
    # Bold keywords use ACCENT_2 theme color (green), normal text is #FFFFFF
    # Font is Akkurat LL throughout
    for i, step in enumerate(next_steps):
        # Add empty paragraph first time to match original P0 being empty
        if i == 0:
            p_empty = tf.paragraphs[0]  # Use existing empty P0
            p = tf.add_paragraph()
        else:
            p_empty = tf.add_paragraph()  # Empty spacer paragraph
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.line_spacing = 1.0
        p.alignment = PP_ALIGN.LEFT

        segments = _parse_rich_segments(step)
        for seg_text, style in segments:
            r = p.add_run()
            r.text = seg_text
            r.font.size = Pt(28)
            r.font.name = "Akkurat LL"
            if style == 'bold':
                r.font.bold = True
                _set_run_scheme_color(r, "accent2")  # Green accent for bold
            elif style == 'green':
                r.font.bold = False
                r.font.color.rgb = OT_GREEN_HIGHLIGHT
            else:
                r.font.bold = False
                r.font.color.rgb = OT_WHITE  # Normal text — white on dark bg

    set_footer_textbox(slide, slide_count(prs))
    if verbose:
        logger.info("Created What's Next slide")


def make_thanks(prs, ot_team, client_contacts, client_name, verbose=False):
    """Thanks slide — 'Bullet Points Ash' layout.

    Original formatting:
    - PICTURE at (0, 0) 5.39x5.99 — team photo full-bleed left
    - TEXT_BOX at (5.91, 1.48) 8.21x1.38 — "Many thanks!" 24pt #FFFFFF left-aligned
    - TEXT_BOX at (6.84, 2.89) 6.36x1.78 — Team names 20pt #FFFFFF, anchor=t
    - Multiple small logo PICTURES at bottom
    """
    layout = find_layout(prs, "Bullet Points Ash")
    if not layout:
        return
    slide = prs.slides.add_slide(layout)

    # Clear all default placeholders from this layout
    for idx in [0, 1, 27, 36, 42, 44, 45, 46, 47]:
        ph = ph_by_idx(slide, idx)
        if ph:
            try:
                sp = ph._element
                sp.getparent().remove(sp)
            except Exception:
                pass

    # Bottom dark rectangle bar (like original Rechteck 27)
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.0), Inches(6.0), Inches(13.33), Inches(1.5)
    )
    rect.line.fill.background()
    # Dark fill matching ash bg
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x24, 0x24, 0x24)  # dk2 ash

    # Left: team photo placeholder (full bleed)
    image_placeholder(slide, 0, 0, 5.39, 5.99, "Add team group photo")

    # "Many thanks!" — 24pt white, Wavetable font
    tb_thanks = slide.shapes.add_textbox(
        Inches(5.91), Inches(1.48), Inches(8.21), Inches(1.38)
    )
    tf_t = tb_thanks.text_frame
    tf_t.word_wrap = True
    tf_t.vertical_anchor = MSO_ANCHOR.TOP
    p = tf_t.paragraphs[0]
    p.line_spacing = 1.0
    r = p.add_run()
    r.text = "Many thanks!"
    r.font.size = Pt(24)
    r.font.color.rgb = OT_WHITE
    r.font.name = "Wavetable"

    # Team names — 20pt white, Akkurat LL font
    team_tb = slide.shapes.add_textbox(
        Inches(6.84), Inches(2.89), Inches(6.36), Inches(1.78)
    )
    tf = team_tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    all_members = ot_team[:]
    if client_contacts:
        all_members.extend(client_contacts)

    for i, member in enumerate(all_members):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = member
        r.font.size = Pt(20)
        r.font.color.rgb = OT_WHITE
        r.font.name = "Akkurat LL"

    # Bottom logos
    image_placeholder(slide, 0.5, 6.5, 2.5, 0.8, "Client logo")
    image_placeholder(slide, 9.6, 6.5, 2.7, 0.8, "OT + partner logos")

    set_footer_textbox(slide, slide_count(prs))
    if verbose:
        logger.info("Created Thanks slide")


# ---------------------------------------------------------------------------
# Main generation
# ---------------------------------------------------------------------------

def generate_presentation(template_path, variables_path, content_path,
                          output_path, verbose=False):
    variables = load_json(variables_path)
    content = load_json(content_path)
    prs = Presentation(str(template_path))

    client = variables.get("client_name", "Client")
    location = variables.get("location", "Location")
    d1 = variables.get("hackathon_dates", {}).get("day1", "DD.MM.YYYY")
    d2 = variables.get("hackathon_dates", {}).get("day2", "DD.MM.YYYY")
    use_case = variables.get("use_case_title", "Use Case")
    ot_team = variables.get("team_members", {}).get("ot_team", [])
    client_contacts = variables.get("team_members", {}).get("client_contacts", [])

    logger.info(f"Generating presentation for {client}")

    # === DAY 1 PRE-INTRO ===
    make_cover(prs, client, location, d1, use_case, day=1, verbose=verbose)
    make_checkin(prs, content.get("check_in", {}).get("questions", []),
                 client, verbose=verbose)
    make_agenda(prs, content.get("agenda", {}), verbose=verbose)

    # === USE CASE SECTION ===
    make_toc(prs, verbose=verbose)

    uc = content.get("use_case", {})

    # 01 Pain — with <<green>> inline highlights
    make_pain(prs, uc.get("pain_points", []), verbose=verbose)

    # 02 Data — with **bold** prefixes
    data_sources = uc.get("data_sources", [])
    data_bullets = [
        f"{src.get('title', '')} {src.get('description', '')}"
        for src in data_sources
    ]
    make_data(prs, data_bullets, verbose=verbose)

    make_data_screenshots(prs, verbose=verbose)

    # 03 Approach — with **bold** step prefixes
    approach = uc.get("approach_steps", [])
    q = uc.get("approach_question", "")
    if q:
        approach = approach + [q]
    make_approach(prs, approach, verbose=verbose)

    # 04 Challenges — with **bold** — format
    make_challenges(prs, uc.get("challenges", []), verbose=verbose)

    # === BREAKTHROUGH DIVIDER ===
    make_divider(prs, "Let's create A BREAKTHROUGH!", verbose=verbose)

    # === DAY 2 ===
    make_cover(prs, client, location, d2, use_case, day=2, verbose=verbose)
    make_image_slide_blank(prs, "Add team photos from Day 1 here", verbose=verbose)

    make_divider(prs, "What have we done in the past 30h?", verbose=verbose)

    # Process Flow & Architecture — separate functions with exact formatting
    make_process_flow(prs, verbose=verbose)
    make_architecture(prs, verbose=verbose)

    # === RESULTS ===
    res = content.get("results", {})
    make_business_value(prs, res.get("business_value", {}), verbose=verbose)

    poc = res.get("poc_summary", {})
    make_poc_summary(prs,
                     poc.get("intro", "We have built a proof of concept (PoC):"),
                     poc.get("features", []),
                     verbose=verbose)

    make_divider(prs, "DEMO", verbose=verbose)
    make_divider(prs, "Expectations check", verbose=verbose)

    # What's Next — dedicated function with manual textboxes
    make_whats_next(prs, res.get("next_steps", []), verbose=verbose)

    make_thanks(prs, ot_team, client_contacts, client, verbose=verbose)

    prs.save(str(output_path))
    logger.info(f"Saved {slide_count(prs)} slides → {output_path}")


def main():
    parser = argparse.ArgumentParser(description="Generate OT hackathon presentation v2")
    parser.add_argument("--template", type=Path, required=True)
    parser.add_argument("--variables", type=Path, required=True)
    parser.add_argument("--content", type=Path, required=True)
    parser.add_argument("--output", type=Path, required=True)
    parser.add_argument("--verbose", action="store_true")
    args = parser.parse_args()

    if args.verbose:
        logger.setLevel(logging.DEBUG)

    try:
        generate_presentation(args.template, args.variables, args.content,
                              args.output, args.verbose)
    except Exception as e:
        logger.error(f"Error: {e}")
        if args.verbose:
            import traceback
            traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()

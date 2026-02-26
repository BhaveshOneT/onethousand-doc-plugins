#!/usr/bin/env python3
"""
Kick-Off Presentation Generator -- One Thousand GmbH

Generates branded kick-off presentations from a template PPTX + JSON inputs.
Fills 20 slides covering: Cover, Agenda (with section dividers), Check-In,
Pain x Data, Hackathon Validation, Step by Step, Architecture, Sprint Goals,
Timeline/Gantt, Progress/Risks, Participants, Meetings, Discussion, Thank You.

Key design decisions:
- Template-first: opens existing 20-slide template and fills placeholders
- Theme colors (no explicit RGB) for most text -- inherits from slide master
- Only #00B050 green highlights and #19A960 accent green are explicit
- Manual TEXT_BOXes where needed for flexible multi-line content
- Agenda divider slides highlight the active section in green, rest in gray

Usage:
    python generate_kickoff_pptx.py \\
        --template assets/templates/ot-kickoff-template.pptx \\
        --variables variables.json \\
        --content content.json \\
        --output output.pptx \\
        [--verbose]
"""

import argparse
import json
import logging
import os
import re
import sys
from copy import deepcopy
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Brand Colors (from OT PPTX theme analysis)
# ---------------------------------------------------------------------------
# Theme mapping: dk1=#000000, lt1=#FFFFFF, dk2=#242424 (ash), lt2=#D5F89E (lime)
# accent2=#19A960 (green)
OT_GREEN_HIGHLIGHT = RGBColor(0x00, 0xB0, 0x50)    # #00B050 -- green inline highlights
OT_GREEN_ACCENT    = RGBColor(0x19, 0xA9, 0x60)     # #19A960 -- scheme accent2
OT_GREEN_SHARP     = RGBColor(0x18, 0xA0, 0x5A)     # #18A05A -- active agenda section
OT_WHITE           = RGBColor(0xFF, 0xFF, 0xFF)      # Explicit white
OT_BLACK           = RGBColor(0x00, 0x00, 0x00)      # Black text
OT_GRAY            = RGBColor(0xBB, 0xBB, 0xBB)      # Inactive agenda items / hint text
OT_ASH             = RGBColor(0x2F, 0x2F, 0x2F)      # Body text on lime bg
OT_LIME_BG         = RGBColor(0xD5, 0xF8, 0x9E)      # Lime background (reference only)
OT_LIME_TEXT       = RGBColor(0xD5, 0xF8, 0x9E)      # Text on dark backgrounds (matches lt2)

# Gantt bar colors
GANTT_GREEN        = RGBColor(0x19, 0xA9, 0x60)      # Sprint bar fill
GANTT_GRAY_BG      = RGBColor(0xF2, 0xF2, 0xF2)      # Empty cell background

logging.basicConfig(level=logging.INFO, format="%(levelname)s - %(message)s")
logger = logging.getLogger(__name__)

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5

# ---------------------------------------------------------------------------
# Rich Text Parsing -- **bold** and <<green>>
# ---------------------------------------------------------------------------
_RICH_RE = re.compile(r"(\*\*.*?\*\*|<<.*?>>)")


def _parse_rich_segments(text: str) -> List[Tuple[str, str]]:
    """Parse text with **bold** and <<green>> markup into segments.

    Returns list of (text, style) tuples where style is 'bold', 'green', or 'normal'.
    """
    parts = _RICH_RE.split(text)
    segments = []
    for part in parts:
        if not part:
            continue
        if part.startswith("**") and part.endswith("**"):
            segments.append((part[2:-2], "bold"))
        elif part.startswith("<<") and part.endswith(">>"):
            segments.append((part[2:-2], "green"))
        else:
            segments.append((part, "normal"))
    return segments


def add_rich_paragraph(
    text_frame,
    text: str,
    font_size: int = 13,
    use_theme_color: bool = True,
    explicit_color: Optional[RGBColor] = None,
    green_color: RGBColor = OT_GREEN_HIGHLIGHT,
    spacing_before: int = 6,
    spacing_after: int = 6,
    is_first: bool = False,
    alignment=None,
    bold_base: bool = False,
    font_name: Optional[str] = None,
):
    """Add a paragraph with mixed bold/green/normal formatting.

    Markup:
      **text** -> bold, theme or explicit color
      <<text>> -> normal, green_color (#00B050)
      plain    -> normal, theme or explicit color

    If use_theme_color=True, normal/bold runs do NOT set font.color.rgb
    (they inherit from theme). Only <<green>> runs get explicit color.
    If explicit_color is set, it overrides theme for normal/bold runs.
    """
    p = text_frame.paragraphs[0] if is_first else text_frame.add_paragraph()
    p.space_before = Pt(spacing_before)
    p.space_after = Pt(spacing_after)
    if alignment:
        p.alignment = alignment

    segments = _parse_rich_segments(text)
    for seg_text, style in segments:
        r = p.add_run()
        r.text = seg_text
        r.font.size = Pt(font_size)
        if font_name:
            r.font.name = font_name
        if style == "bold":
            r.font.bold = True
            if explicit_color:
                r.font.color.rgb = explicit_color
        elif style == "green":
            r.font.bold = False
            r.font.color.rgb = green_color
        else:
            r.font.bold = bold_base
            if explicit_color:
                r.font.color.rgb = explicit_color
    return p


# ---------------------------------------------------------------------------
# Utility Helpers
# ---------------------------------------------------------------------------


def load_json(path: Path) -> Dict[str, Any]:
    """Load and return a JSON file as a dictionary."""
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def find_layout(prs, name: str):
    """Find a slide layout by name. Returns None if not found."""
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    logger.warning(f"Layout not found: {name}")
    return None


def _find_placeholder(slide, idx: int):
    """Find a placeholder on a slide by its idx. Returns None if not found."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _set_placeholder_text(
    slide,
    idx: int,
    text: str,
    size: Optional[int] = None,
    bold: bool = False,
    color: Optional[RGBColor] = None,
    font_name: Optional[str] = None,
):
    """Set text on a placeholder by idx. Logs warning if placeholder not found.

    If color is None, text inherits theme color (no explicit RGB set).
    """
    ph = _find_placeholder(slide, idx)
    if ph is None:
        logger.warning(f"Placeholder idx={idx} not found on slide")
        return
    if not ph.has_text_frame:
        logger.warning(f"Placeholder idx={idx} has no text frame")
        return
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    if size:
        r.font.size = Pt(size)
    if bold:
        r.font.bold = True
    if color:
        r.font.color.rgb = color
    if font_name:
        r.font.name = font_name


def _set_placeholder_text_theme(
    slide,
    idx: int,
    text: str,
    size: Optional[int] = None,
    bold: bool = False,
    font_name: Optional[str] = None,
):
    """Set text on a placeholder using THEME color (no explicit RGB).

    This matches the original PPTX behavior where most text inherits from theme.
    """
    _set_placeholder_text(slide, idx, text, size=size, bold=bold, color=None,
                          font_name=font_name)


def _add_image_to_slide(
    slide,
    image_path: str,
    left: float,
    top: float,
    width: float,
    height: float,
) -> bool:
    """Add an image to the slide at the given position (in inches).

    Returns True if image was added, False if path is invalid or file not found.
    """
    if not image_path:
        logger.info("No image path provided, skipping")
        return False
    if not os.path.isfile(image_path):
        logger.warning(f"Image file not found: {image_path}")
        return False
    try:
        slide.shapes.add_picture(
            image_path,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        logger.debug(f"Added image: {image_path}")
        return True
    except Exception as e:
        logger.warning(f"Failed to add image {image_path}: {e}")
        return False


def reposition_shape(shape, left=None, top=None, width=None, height=None):
    """Reposition/resize a shape using inches."""
    if left is not None:
        shape.left = Inches(left)
    if top is not None:
        shape.top = Inches(top)
    if width is not None:
        shape.width = Inches(width)
    if height is not None:
        shape.height = Inches(height)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=14,
    text_color=None,
    alignment=PP_ALIGN.LEFT,
    bold=False,
    word_wrap=True,
    vertical_anchor=MSO_ANCHOR.TOP,
    font_name=None,
):
    """Add a textbox with optional explicit color."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = vertical_anchor
    p = tf.paragraphs[0]
    p.alignment = alignment
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    if text_color:
        r.font.color.rgb = text_color
    if font_name:
        r.font.name = font_name
    return tb


def add_textbox_theme(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=14,
    alignment=PP_ALIGN.LEFT,
    bold=False,
    word_wrap=True,
    vertical_anchor=MSO_ANCHOR.TOP,
    font_name=None,
):
    """Add a textbox using THEME color (no explicit RGB set on font)."""
    return add_textbox(
        slide, left, top, width, height, text,
        font_size=font_size, text_color=None, alignment=alignment,
        bold=bold, word_wrap=word_wrap, vertical_anchor=vertical_anchor,
        font_name=font_name,
    )


def _cover_logo_placeholder(slide, left, top, width, height):
    """Add a client logo placeholder with gray background fill (#E0E0E0), centered text."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "[IMAGE: Client logo]"
    r.font.size = Pt(14)
    r.font.color.rgb = OT_GRAY
    # Apply gray background fill
    fill = tb.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)


def image_placeholder(slide, left, top, width, height, description):
    """Add a subtle image placeholder with hint text (no heavy bg fill)."""
    add_textbox(
        slide, left, top, width, height,
        f"[IMAGE: {description}]",
        font_size=14,
        text_color=OT_GRAY,
        alignment=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


def add_speaker_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def _ensure_bullet(paragraph, indent=Inches(0.25), margin_left=Inches(0.1)):
    """Ensure a paragraph has bullet-character formatting (•) with proper indent.

    Adds <a:buChar char="•"/> to the paragraph properties if no bullet
    definition already exists. Also sets indent and margin so the bullet
    aligns consistently across all paragraphs (not just the first one which
    inherits template formatting from tf.clear()).
    """
    pPr = paragraph._p.get_or_add_pPr()
    # Check if any bullet definition already exists
    for child in pPr:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag.startswith("bu") and tag != "buNone":
            return  # Already has a bullet definition
    # Remove any buNone
    for child in list(pPr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag == "buNone":
            pPr.remove(child)
    # Set indent and margin for proper bullet alignment
    pPr.set("indent", str(-indent))
    pPr.set("marL", str(margin_left + indent))
    # Add bullet character
    etree.SubElement(pPr, f"{{{NS_A}}}buChar", attrib={"char": "\u2022"})


def _enable_autofit(text_frame):
    """Enable auto-shrink text fitting via <a:normAutofit/>."""
    body_props = text_frame._txBody.find(f"{{{NS_A}}}bodyPr")
    if body_props is None:
        return
    for child in list(body_props):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("noAutofit", "normAutofit", "spAutoFit"):
            body_props.remove(child)
    etree.SubElement(body_props, f"{{{NS_A}}}normAutofit")


def _promote_layout_fill_to_slide(slide_ph, layout):
    """Copy fill and geometry from the layout placeholder to the slide placeholder.

    LibreOffice renders layout fills at the *layout* position, not the slide
    position.  When the slide repositions a placeholder, this causes a mismatch.
    Copying the fill to the slide level ensures it renders at the slide position
    in both PowerPoint and LibreOffice.
    """
    ph_idx = slide_ph.placeholder_format.idx
    # Find matching layout placeholder
    layout_ph = None
    for lph in layout.placeholders:
        if lph.placeholder_format.idx == ph_idx:
            layout_ph = lph
            break
    if layout_ph is None:
        return

    layout_spPr = layout_ph._element.find(f"{{{NS_P}}}spPr")
    slide_spPr = slide_ph._element.find(f"{{{NS_P}}}spPr")
    if layout_spPr is None or slide_spPr is None:
        return

    # Copy fill from layout → slide (replace any existing fill/noFill)
    for child in list(slide_spPr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("solidFill", "gradFill", "pattFill", "blipFill", "noFill"):
            slide_spPr.remove(child)

    layout_fill = None
    for child in layout_spPr:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("solidFill", "gradFill", "pattFill", "blipFill"):
            layout_fill = child
            break
    if layout_fill is not None:
        slide_spPr.append(deepcopy(layout_fill))

    # Copy geometry (prstGeom) if slide doesn't have it
    slide_geom = slide_spPr.find(f"{{{NS_A}}}prstGeom")
    if slide_geom is None:
        layout_geom = layout_spPr.find(f"{{{NS_A}}}prstGeom")
        if layout_geom is not None:
            slide_spPr.append(deepcopy(layout_geom))

    # Copy line from layout
    for child in list(slide_spPr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag == "ln":
            slide_spPr.remove(child)
    layout_ln = layout_spPr.find(f"{{{NS_A}}}ln")
    if layout_ln is not None:
        slide_spPr.append(deepcopy(layout_ln))

    # Remove the fill from the LAYOUT placeholder so LibreOffice doesn't
    # render it at the (different) layout position as a ghost background.
    if layout_fill is not None:
        layout_spPr.remove(layout_fill)
        etree.SubElement(layout_spPr, f"{{{NS_A}}}noFill")
    # Also suppress the layout line to avoid ghost outlines
    layout_ln_elem = layout_spPr.find(f"{{{NS_A}}}ln")
    if layout_ln_elem is not None:
        layout_spPr.remove(layout_ln_elem)


def remove_shape_by_ph_idx(slide, idx):
    """Remove a placeholder element entirely from the slide XML."""
    ph = _find_placeholder(slide, idx)
    if ph is not None:
        sp = ph._element
        sp.getparent().remove(sp)


def _set_run_scheme_color(run, scheme_val="tx1"):
    """Set a run's font color using scheme color reference.

    Common values: tx1=dark text, bg2=lime, accent2=green
    """
    rPr = run._r.get_or_add_rPr()
    for child in list(rPr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag == "solidFill":
            rPr.remove(child)
    solidFill = etree.SubElement(rPr, f"{{{NS_A}}}solidFill")
    schemeClr = etree.SubElement(solidFill, f"{{{NS_A}}}schemeClr")
    schemeClr.set("val", scheme_val)


def _set_cell_fill(cell, color: RGBColor):
    """Set the fill color of a table cell."""
    cell_elem = cell._tc
    tcPr = cell_elem.find(f"{{{NS_A}}}tcPr")
    if tcPr is None:
        tcPr = etree.SubElement(cell_elem, f"{{{NS_A}}}tcPr")
    # Remove existing fill
    for child in list(tcPr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("solidFill", "noFill"):
            tcPr.remove(child)
    solidFill = etree.SubElement(tcPr, f"{{{NS_A}}}solidFill")
    srgbClr = etree.SubElement(solidFill, f"{{{NS_A}}}srgbClr")
    srgbClr.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")


def _set_cell_text(cell, text: str, size: int = 8, bold: bool = False,
                   color: Optional[RGBColor] = None, alignment=PP_ALIGN.CENTER):
    """Set text in a table cell."""
    tf = cell.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    if color:
        r.font.color.rgb = color


# ---------------------------------------------------------------------------
# Agenda / ToC Helpers
# ---------------------------------------------------------------------------

# Mapping of section index (1-6) to placeholder indices on ToC layout
SECTION_PH_MAP = {
    1: {"name": 13, "number": 14, "page": 17},
    2: {"name": 15, "number": 16, "page": 18},
    3: {"name": 19, "number": 20, "page": 21},
    4: {"name": 25, "number": 26, "page": 29},
    5: {"name": 27, "number": 28, "page": 30},
    6: {"name": 31, "number": 32, "page": 33},
}

# Which agenda section is highlighted on each divider slide
# slide_index -> active section number (0 = no highlight = full agenda)
DIVIDER_HIGHLIGHT_MAP = {
    1: 0,   # Full agenda -- no highlight
    2: 1,   # Check In
    4: 2,   # Use-case
    9: 3,   # Timeline
    13: 4,  # Collaboration
    16: 5,  # Other Topics
    18: 6,  # Check Out
}


def _compute_page_numbers(total_slides: int, extra_uc_slides: int = 0) -> Dict[int, str]:
    """Compute page references for each section based on slide positions.

    Base slide positions (0-indexed):
      Section 1 (Check In):      slide 2  (divider), content starts slide 3
      Section 2 (Use-case):      slide 4  (divider), content starts slide 5
      Section 3 (Timeline):      slide 9  (divider) + extra_uc_slides offset
      Section 4 (Collaboration): slide 13 + extra_uc_slides offset
      Section 5 (Other Topics):  slide 16 + extra_uc_slides offset
      Section 6 (Check Out):     slide 18 + extra_uc_slides offset

    Returns dict mapping section number to two-digit page string.
    Page numbers shown in the ToC represent the divider slide position (1-indexed).
    """
    offset = extra_uc_slides
    return {
        1: f"{3:02d}",
        2: f"{5:02d}",
        3: f"{10 + offset:02d}",
        4: f"{14 + offset:02d}",
        5: f"{17 + offset:02d}",
        6: f"{19 + offset:02d}",
    }


def _fill_agenda_slide(
    slide,
    sections: List[Dict[str, str]],
    page_numbers: Dict[int, str],
    active_section: int = 0,
):
    """Fill a ToC / agenda divider slide with section names, numbers, and pages.

    Args:
        slide: The slide object (ToC middle w/o lines layout).
        sections: List of dicts with 'name' and 'number' keys.
        page_numbers: Dict mapping section number (1-6) to page string.
        active_section: Which section to highlight (1-6), or 0 for no highlight.
    """
    for i, section in enumerate(sections[:6], start=1):
        ph_info = SECTION_PH_MAP.get(i)
        if not ph_info:
            continue

        name = section.get("name", f"Section {i}")
        number = section.get("number", f"{i:02d}")
        page = page_numbers.get(i, "")

        # Determine color: active = green, inactive = gray, full agenda = theme
        if active_section == 0:
            # Full agenda -- all sections in theme color (no explicit color)
            name_color = None
            num_color = None
            page_color = None
        elif i == active_section:
            # Active section -- sharp green
            name_color = OT_GREEN_SHARP
            num_color = OT_GREEN_SHARP
            page_color = OT_GREEN_SHARP
        else:
            # Inactive section -- gray
            name_color = OT_GRAY
            num_color = OT_GRAY
            page_color = OT_GRAY

        _set_placeholder_text(slide, ph_info["name"], name, color=name_color)
        _set_placeholder_text(slide, ph_info["number"], number, color=num_color)
        _set_placeholder_text(slide, ph_info["page"], page, color=page_color)


# ---------------------------------------------------------------------------
# Slide Population Functions
# ---------------------------------------------------------------------------


def fill_cover(slide, variables: Dict, content: Dict, verbose: bool = False):
    """Fill cover slide (slide 0) -- Title Lime + one Logo layout.

    Elements:
    - idx=10: Date
    - idx=0: Title (tagline + use-case subtitle)
    - Client logo image (if provided)
    """
    date = variables.get("date", "DD.MM.YYYY")
    tagline = variables.get("tagline", "")
    project_title = variables.get("project_title", "")
    use_cases = variables.get("use_cases", [])
    images = variables.get("images", {})

    # Date placeholder -- format: "Project Kick-Off | DD.MM.YYYY"
    language = variables.get("language", "en")
    kickoff_label = "Projekt Kick-Off" if language == "de" else "Project Kick-Off"
    date_text = f"{kickoff_label} | {date}"
    _set_placeholder_text(slide, 10, date_text, size=13, font_name="Akkurat LL",
                          color=OT_BLACK)

    # Title -- tagline on line 1, use-case subtitle on line 2
    title_ph = _find_placeholder(slide, 0)
    if title_ph and title_ph.has_text_frame:
        tf = title_ph.text_frame
        tf.clear()

        # Paragraph 1: Tagline (bold, 32pt)
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = tagline if tagline else project_title
        r1.font.size = Pt(32)
        r1.font.bold = True

        # Paragraph 2: Use-case subtitle (not bold, 16pt)
        if use_cases:
            uc_title = use_cases[0].get("title", "")
            if uc_title:
                p2 = tf.add_paragraph()
                r2 = p2.add_run()
                r2.text = f"Kick-Off | {uc_title}"
                r2.font.size = Pt(16)
                r2.font.bold = False

        _enable_autofit(tf)
    elif title_ph is None:
        logger.warning("Cover: title placeholder idx=0 not found")

    # Client logo -- positioned to match reference: (10.07, 4.99) 3.09x3.09
    logo_path = images.get("client_logo", "")
    logo_left, logo_top, logo_w, logo_h = 10.07, 4.99, 3.09, 3.09
    if logo_path:
        if not _add_image_to_slide(slide, logo_path, logo_left, logo_top, logo_w, logo_h):
            _cover_logo_placeholder(slide, logo_left, logo_top, logo_w, logo_h)
    else:
        _cover_logo_placeholder(slide, logo_left, logo_top, logo_w, logo_h)

    if verbose:
        logger.info("Filled cover slide")


def fill_agenda(
    slide,
    sections: List[Dict[str, str]],
    page_numbers: Dict[int, str],
    active_section: int,
    slide_index: int,
    verbose: bool = False,
):
    """Fill an agenda / ToC divider slide.

    Args:
        slide: Slide object.
        sections: List of section dicts with 'name' and 'number'.
        page_numbers: Dict mapping section (1-6) to page string.
        active_section: 0 for full agenda, 1-6 for divider highlighting.
        slide_index: Current slide index for logging.
    """
    _fill_agenda_slide(slide, sections, page_numbers, active_section)

    label = "Full Agenda" if active_section == 0 else f"Agenda Divider (section {active_section})"
    if verbose:
        logger.info(f"Filled slide {slide_index}: {label}")


def fill_checkin(slide, content: Dict, variables: Dict, images: Dict,
                 verbose: bool = False):
    """Fill check-in slide (slide 3) -- Bullet Points Lime layout.

    Elements:
    - idx=0: Title ("CHECK IN" or custom)
    - Existing textbox id=6 at (~0.38, ~3.09) for questions
    - Optional GIF image on the right
    """
    check_in = content.get("check_in", {})
    title = check_in.get("title", "CHECK IN")
    questions = check_in.get("questions", [])

    _set_placeholder_text(slide, 0, title, size=20, bold=True)
    title_ph = _find_placeholder(slide, 0)
    if title_ph and title_ph.has_text_frame:
        _enable_autofit(title_ph.text_frame)

    # Clear subtitle
    sub_ph = _find_placeholder(slide, 1)
    if sub_ph:
        _set_placeholder_text_theme(slide, 1, "")

    # Find the existing question textbox by position (id=6, near 0.38, 3.09)
    question_box = None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            if shape.placeholder_format is not None:
                continue
        except:
            pass
        l = (shape.left or 0) / 914400
        t = (shape.top or 0) / 914400
        w = (shape.width or 0) / 914400
        h = (shape.height or 0) / 914400
        # Match the template textbox: near left edge, mid-slide vertically, ~6" wide
        if 0.0 < l < 1.0 and 2.5 < t < 4.0 and w > 4.0 and h > 0.8:
            question_box = shape
            break

    if question_box and questions:
        client_name = variables.get("client_name", "Client")
        tf = question_box.text_frame
        tf.clear()
        tf.word_wrap = True
        _enable_autofit(tf)

        for i, q in enumerate(questions):
            q = q.replace("{client_name}", client_name)
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(8)
            p.space_after = Pt(8)
            # Ensure bullet formatting on all paragraphs (not just the first)
            _ensure_bullet(p)
            r = p.add_run()
            r.text = q
            r.font.size = Pt(18)
            r.font.color.rgb = OT_ASH
    elif questions:
        # Fallback: create textbox if template one not found
        client_name = variables.get("client_name", "Client")
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(3.0), Inches(6.0), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        _enable_autofit(tf)
        for i, q in enumerate(questions):
            q = q.replace("{client_name}", client_name)
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(8)
            p.space_after = Pt(8)
            _ensure_bullet(p)
            r = p.add_run()
            r.text = q
            r.font.size = Pt(18)
            r.font.color.rgb = OT_ASH

    # Optional GIF -- reference position: (6.79, 2.01) 6.17x3.47
    gif_left, gif_top, gif_w, gif_h = 6.79, 2.01, 6.17, 3.47
    gif_path = images.get("checkin_gif", "")
    if gif_path:
        if not _add_image_to_slide(slide, gif_path, gif_left, gif_top, gif_w, gif_h):
            _checkin_gif_placeholder(slide, gif_left, gif_top, gif_w, gif_h)
    else:
        _checkin_gif_placeholder(slide, gif_left, gif_top, gif_w, gif_h)

    if verbose:
        logger.info("Filled check-in slide")


def _checkin_gif_placeholder(slide, left, top, width, height):
    """Add a check-in GIF placeholder with gray background (#E0E0E0), centered text."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "[IMAGE: Add icebreaker GIF]"
    r.font.size = Pt(14)
    r.font.color.rgb = OT_GRAY
    # Apply gray background fill
    fill = tb.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)


def fill_pain_data(
    slide,
    uc_data: Dict,
    images: Dict,
    verbose: bool = False,
):
    """Fill Pain x Data slide (slide 5) -- Bullet Points Lime layout.

    Elements:
    - idx=0: Title (UC title or "PAIN x DATA")
    - Existing template text boxes for pain, data, solution bullets
    - Optional PxD GIF
    """
    title = uc_data.get("title", "PAIN x DATA")
    pain_points = uc_data.get("pain_points", [])
    data_sources = uc_data.get("data_sources", [])
    solution = uc_data.get("solution", [])

    _set_placeholder_text(slide, 0, title.upper(), size=20, bold=True)
    title_ph = _find_placeholder(slide, 0)
    if title_ph and title_ph.has_text_frame:
        _enable_autofit(title_ph.text_frame)

    # Clear subtitle
    sub_ph = _find_placeholder(slide, 1)
    if sub_ph:
        _set_placeholder_text_theme(slide, 1, "")

    # Find existing template textboxes by position
    # Template has: pain box ~(0.25, 4.72), data box ~(4.97, 4.72), solution box ~(9.37, 4.71)
    def _is_placeholder(shape):
        try:
            return shape.placeholder_format is not None
        except Exception:
            return False

    def _find_content_box(shapes, target_left, target_top, tolerance=1.0):
        best = None
        best_dist = float('inf')
        for shape in shapes:
            if not shape.has_text_frame:
                continue
            if _is_placeholder(shape):
                continue
            l = (shape.left or 0) / 914400
            t = (shape.top or 0) / 914400
            w = (shape.width or 0) / 914400
            h = (shape.height or 0) / 914400
            # Only consider reasonably sized text boxes (not the giant P/D watermarks or tiny copyright)
            if w < 1.5 or h < 1.0 or w > 5.0:
                continue
            dist = abs(l - target_left) + abs(t - target_top)
            if dist < best_dist and dist < tolerance:
                best_dist = dist
                best = shape
        return best

    pain_box = _find_content_box(slide.shapes, 0.25, 4.72)
    data_box = _find_content_box(slide.shapes, 4.97, 4.72)
    solution_box = _find_content_box(slide.shapes, 9.37, 4.71)

    # Fill pain box
    if pain_box and pain_points:
        tf = pain_box.text_frame
        tf.clear()
        tf.word_wrap = True
        _enable_autofit(tf)
        for i, bullet in enumerate(pain_points):
            add_rich_paragraph(
                tf, bullet, font_size=10, use_theme_color=False,
                explicit_color=OT_ASH, green_color=OT_GREEN_HIGHLIGHT,
                spacing_before=2, spacing_after=2, is_first=(i == 0),
            )

    # Fill data box
    if data_box and data_sources:
        tf = data_box.text_frame
        tf.clear()
        tf.word_wrap = True
        _enable_autofit(tf)
        for i, bullet in enumerate(data_sources):
            add_rich_paragraph(
                tf, bullet, font_size=10, use_theme_color=False,
                explicit_color=OT_ASH, green_color=OT_GREEN_HIGHLIGHT,
                spacing_before=2, spacing_after=2, is_first=(i == 0),
            )

    # Fill solution box
    if solution_box and solution:
        tf = solution_box.text_frame
        tf.clear()
        tf.word_wrap = True
        _enable_autofit(tf)
        for i, bullet in enumerate(solution):
            add_rich_paragraph(
                tf, bullet, font_size=10, use_theme_color=False,
                explicit_color=OT_ASH, green_color=OT_GREEN_HIGHLIGHT,
                spacing_before=2, spacing_after=2, is_first=(i == 0),
            )

    # Optional PxD GIF -- only add if file exists
    gif_path = images.get("pxd_gif", "")
    if gif_path and os.path.isfile(gif_path):
        _add_image_to_slide(slide, gif_path, 8.5, 1.8, 4.2, 4.5)

    if verbose:
        logger.info("Filled Pain x Data slide")


def fill_hackathon(slide, content: Dict, images: Dict, verbose: bool = False):
    """Fill hackathon validation slide (slide 6) -- DEFAULT layout.

    Elements:
    - TEXT_BOX for title
    - TEXT_BOX for highlight bullets
    - Optional hackathon images
    """
    hackathon = content.get("hackathon", {})
    title = hackathon.get("title", "HACKATHON VALIDATION")
    highlights = hackathon.get("highlights", [])

    # Title textbox (DEFAULT layout has no usable title placeholder)
    add_textbox(
        slide, 0.5, 0.5, 12.0, 1.0,
        title.upper(),
        font_size=24, text_color=OT_WHITE, bold=False,
        alignment=PP_ALIGN.LEFT,
        vertical_anchor=MSO_ANCHOR.TOP,
    )

    # Highlight bullets
    if highlights:
        tb = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.8), Inches(5.5), Inches(4.0)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        _enable_autofit(tf)

        for i, h in enumerate(highlights):
            add_rich_paragraph(
                tf, h,
                font_size=16,
                use_theme_color=False,
                explicit_color=OT_WHITE,
                green_color=OT_GREEN_HIGHLIGHT,
                spacing_before=6, spacing_after=6,
                is_first=(i == 0),
            )

    # Hackathon images
    hack_images = images.get("hackathon_images", [])
    if hack_images:
        # Place up to 3 images in a row on the right side
        img_left = 6.5
        img_top = 1.8
        img_w = 3.0
        img_h = 2.0
        for j, img_path in enumerate(hack_images[:3]):
            y = img_top + j * (img_h + 0.3)
            if not _add_image_to_slide(slide, img_path, img_left, y, img_w, img_h):
                image_placeholder(slide, img_left, y, img_w, img_h,
                                  f"Hackathon image {j+1}")
    else:
        image_placeholder(slide, 6.5, 2.0, 5.5, 4.0,
                          "Add hackathon / Miro board images")

    if verbose:
        logger.info("Filled hackathon validation slide")


def fill_step_by_step(slide, content: Dict, images: Dict, verbose: bool = False):
    """Fill step-by-step slide (slide 7) -- Bullet Points Lime layout.

    Elements:
    - idx=0: Title
    - 4 existing phase label textboxes along the bottom (~y 5.2-5.5)
    - Image area above the labels (~y 2.0-5.0)
    """
    sbs = content.get("step_by_step", {})
    title = sbs.get("title", "STEP BY STEP")
    phases = sbs.get("phases", [])

    _set_placeholder_text(slide, 0, title.upper(), size=20, bold=True)
    title_ph = _find_placeholder(slide, 0)
    if title_ph and title_ph.has_text_frame:
        _enable_autofit(title_ph.text_frame)

    # Clear subtitle
    sub_ph = _find_placeholder(slide, 1)
    if sub_ph:
        _set_placeholder_text_theme(slide, 1, "")

    # Find the 4 phase label textboxes sorted by left position
    phase_boxes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            if shape.placeholder_format is not None:
                continue
        except:
            pass
        l = (shape.left or 0) / 914400
        t = (shape.top or 0) / 914400
        w = (shape.width or 0) / 914400
        h = (shape.height or 0) / 914400
        # Phase labels are along the bottom (y ~5.2-5.5), 2-3" wide
        if 4.8 < t < 6.0 and w > 1.5 and w < 4.0 and h < 1.0:
            phase_boxes.append((l, shape))

    phase_boxes.sort(key=lambda x: x[0])

    # Fill phase labels (name only -- descriptions are not rendered on this slide)
    for i, phase in enumerate(phases[:4]):
        name = phase.get("name", f"Phase {i+1}")

        if i < len(phase_boxes):
            shape = phase_boxes[i][1]
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            _enable_autofit(tf)

            p = tf.paragraphs[0]
            # Phase name only (bold, green, 14pt)
            r_name = p.add_run()
            r_name.text = name
            r_name.font.size = Pt(14)
            r_name.font.bold = True
            r_name.font.color.rgb = OT_GREEN_ACCENT

    # Screenshots/images area (above the phase labels)
    step_screenshots = images.get("step_screenshots", [])
    if step_screenshots:
        img_y = 2.0
        img_h = 3.0
        img_w = 2.8
        for j, ss_path in enumerate(step_screenshots[:4]):
            x = 0.8 + j * 3.0
            _add_image_to_slide(slide, ss_path, x, img_y, img_w, img_h)
    else:
        image_placeholder(slide, 0.8, 2.0, 11.5, 3.0, "Add step-by-step screenshots")

    if verbose:
        logger.info("Filled step-by-step slide")


def fill_architecture(slide, content: Dict, images: Dict, verbose: bool = False):
    """Fill architecture slide (slide 8) -- Bullet Points Lime layout.

    Elements:
    - idx=0: Title
    - Existing autoshape (rounded rect) at (~9.28, ~2.05, ~3.76x3.77) for capabilities
    - Left area for architecture diagram image
    """
    arch = content.get("architecture", {})
    title = arch.get("title", "ARCHITECTURE")
    capabilities = arch.get("capabilities", [])

    _set_placeholder_text(slide, 0, title.upper(), size=20, bold=True)
    title_ph = _find_placeholder(slide, 0)
    if title_ph and title_ph.has_text_frame:
        _enable_autofit(title_ph.text_frame)

    # Clear subtitle
    sub_ph = _find_placeholder(slide, 1)
    if sub_ph:
        _set_placeholder_text_theme(slide, 1, "")

    # Architecture diagram area (left side)
    diagram_path = images.get("architecture_diagram", "")
    if diagram_path and os.path.isfile(diagram_path):
        _add_image_to_slide(slide, diagram_path, 0.5, 2.0, 8.0, 4.5)
    else:
        image_placeholder(slide, 0.5, 2.0, 8.0, 4.5, "Add architecture diagram")

    # Find the capabilities autoshape (rounded rect at ~9.28, 2.05, ~3.76x3.77)
    cap_shape = None
    for shape in slide.shapes:
        l = (shape.left or 0) / 914400
        t = (shape.top or 0) / 914400
        w = (shape.width or 0) / 914400
        h = (shape.height or 0) / 914400
        if shape.has_text_frame and 8.5 < l < 10.5 and 1.5 < t < 3.0 and w > 2.5 and h > 2.5:
            try:
                if shape.placeholder_format is not None:
                    continue
            except:
                pass
            cap_shape = shape
            break

    if cap_shape and capabilities:
        # Write directly into the autoshape's text frame with generous internal
        # margins.  The top margin (0.5") clears the decorative green curve that
        # crosses the upper portion of the dark card.  Using the autoshape's own
        # text frame (instead of an overlay textbox) avoids z-order rendering
        # differences between PowerPoint and LibreOffice.
        tf = cap_shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_top = Inches(0.5)
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_bottom = Inches(0.1)

        cap_font_size = 7 if len(capabilities) >= 5 else 8
        for i, cap in enumerate(capabilities):
            add_rich_paragraph(
                tf, cap, font_size=cap_font_size,
                use_theme_color=False, explicit_color=OT_LIME_TEXT,
                green_color=OT_GREEN_HIGHLIGHT,
                spacing_before=1, spacing_after=1,
                is_first=(i == 0),
            )
    elif capabilities:
        # Fallback: add textbox if autoshape not found
        tb = slide.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(4.0), Inches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        _enable_autofit(tf)
        cap_font_size = 9 if len(capabilities) <= 5 else 8
        for i, cap in enumerate(capabilities):
            add_rich_paragraph(
                tf, cap, font_size=cap_font_size,
                use_theme_color=False, explicit_color=OT_LIME_TEXT,
                green_color=OT_GREEN_HIGHLIGHT,
                spacing_before=2, spacing_after=2,
                is_first=(i == 0),
            )

    if verbose:
        logger.info("Filled architecture slide")


def fill_sprint_goals(slide, content: Dict, verbose: bool = False):
    """Fill sprint goals slide (slide 10) -- 1_Bullet Points Lime layout.

    Elements:
    - idx=0: CENTER_TITLE
    - idx=1: SUBTITLE
    - idx=36, 46, 48, 47: Sprint title placeholders (y≈2.29")
    - idx=42, 45, 49, 50: Sprint description placeholders (y≈3.02")
    """
    sprints = content.get("sprints", [])
    title_text = "SPRINT GOALS"

    _set_placeholder_text(slide, 0, title_text, size=20, bold=True)
    title_ph0 = _find_placeholder(slide, 0)
    if title_ph0 and title_ph0.has_text_frame:
        _enable_autofit(title_ph0.text_frame)

    # Clear subtitle
    sub_ph = _find_placeholder(slide, 1)
    if sub_ph:
        _set_placeholder_text_theme(slide, 1, "")

    # Sprint placeholders mapping: (title_idx, description_idx)
    # These are placeholder FORMAT indices (not shape IDs)
    # Description text uses explicit color #292A2E (not theme color) on light cards
    SPRINT_DESC_COLOR = RGBColor(0x29, 0x2A, 0x2E)

    # Dark-background sprint card indices (these have solidFill tx2 in the layout)
    DARK_TITLE_INDICES = {46, 47}

    # Dark-background description indices (paired with dark titles)
    DARK_DESC_INDICES = {45, 50}

    sprint_slots = [
        (36, 42),   # Sprint 1 - light card
        (46, 45),   # Sprint 2 - dark card
        (48, 49),   # Sprint 3 - light card
        (47, 50),   # Sprint 4 - dark card
    ]

    layout = slide.slide_layout

    for i, (title_idx, desc_idx) in enumerate(sprint_slots):
        if i < len(sprints):
            sprint = sprints[i]
            s_title = sprint.get("title", f"Sprint {i+1}")
            s_desc = sprint.get("description", "")

            is_dark = title_idx in DARK_TITLE_INDICES
            title_color = OT_LIME_TEXT if is_dark else None
            desc_color = OT_LIME_TEXT if is_dark else SPRINT_DESC_COLOR

            # Sprint title: break "Sprint N: Name" onto two lines for narrow
            # cards (~2.2" wide) to prevent truncation in LibreOffice.
            display_title = s_title
            if ": " in s_title:
                parts = s_title.split(": ", 1)
                display_title = f"{parts[0]}:\n{parts[1]}"

            _set_placeholder_text(slide, title_idx, display_title, size=10, bold=True,
                                  color=title_color)
            s_title_ph = _find_placeholder(slide, title_idx)
            if s_title_ph and s_title_ph.has_text_frame:
                _enable_autofit(s_title_ph.text_frame)

            _set_placeholder_text(slide, desc_idx, s_desc, size=7,
                                  color=desc_color)
            desc_ph = _find_placeholder(slide, desc_idx)
            if desc_ph and desc_ph.has_text_frame:
                _enable_autofit(desc_ph.text_frame)

            # For dark cards, promote the layout fill/geometry to the slide
            # level.  LibreOffice renders layout fills at the *layout* position,
            # not the slide position — so if the slide repositioned the
            # placeholder, the fill and text clip region would be misaligned.
            if is_dark:
                if s_title_ph:
                    _promote_layout_fill_to_slide(s_title_ph, layout)
                if desc_ph:
                    _promote_layout_fill_to_slide(desc_ph, layout)
        else:
            _set_placeholder_text_theme(slide, title_idx, "")
            _set_placeholder_text_theme(slide, desc_idx, "")

    if verbose:
        logger.info("Filled sprint goals slide")


def fill_timeline(slide, content: Dict, verbose: bool = False):
    """Fill timeline/Gantt slide (slide 11) -- Calendar Lime w/o lines layout.

    Elements:
    - idx=0: CENTER_TITLE
    - idx=21: Year/quarter label
    - idx=98, 99: Legend labels
    - TABLE shape: row 0 = months, row 1 = weeks, rows 2+ = tasks with Gantt bars
    """
    timeline = content.get("timeline", {})
    title = timeline.get("title", "TIMELINE")
    year_quarter = timeline.get("year_quarter", "")
    legend = timeline.get("legend", {})
    months = timeline.get("months", [])
    weeks = timeline.get("weeks", [])
    tasks = timeline.get("tasks", [])

    _set_placeholder_text(slide, 0, title.upper(), size=20, bold=True,
                          font_name="Wavetable")
    title_ph = _find_placeholder(slide, 0)
    if title_ph and title_ph.has_text_frame:
        _enable_autofit(title_ph.text_frame)

    # Year/quarter label
    if year_quarter:
        _set_placeholder_text(slide, 21, year_quarter, size=10)

    # Legend labels -- widen placeholders to fit longer client names
    legend_client = legend.get("client", "Client")
    legend_ot = legend.get("ot", "1000")
    _set_placeholder_text(slide, 98, legend_client, size=7)
    _set_placeholder_text(slide, 99, legend_ot, size=7)

    # Widen the client legend box if the name is long (template default is 0.61")
    # and shift the OT legend box (ph=99) right so it doesn't overlap.
    legend_ph = _find_placeholder(slide, 98)
    ot_legend_ph = _find_placeholder(slide, 99)
    if legend_ph:
        min_width = max(Inches(0.6), Inches(0.12 * len(legend_client)))
        if legend_ph.width < min_width:
            legend_ph.width = min_width
        if legend_ph.has_text_frame:
            _enable_autofit(legend_ph.text_frame)
        # Ensure ph=99 starts right after ph=98 with a small gap
        if ot_legend_ph:
            gap = Inches(0.05)
            new_left = legend_ph.left + legend_ph.width + gap
            if new_left > ot_legend_ph.left:
                ot_legend_ph.left = new_left

    # Find the TABLE shape on the slide
    table_shape = None
    for shape in slide.shapes:
        if shape.has_table:
            table_shape = shape
            break

    if table_shape is None:
        logger.warning("Timeline slide: no TABLE shape found. Cannot fill Gantt chart.")
        if verbose:
            logger.info("Filled timeline slide (no table found)")
        return

    table = table_shape.table
    num_rows = len(table.rows)
    num_cols = len(table.columns)

    logger.debug(f"Timeline table: {num_rows} rows x {num_cols} cols")

    # Row 0: Month headers (merged cells spanning multiple columns)
    # Column 0 is the label column; columns 1+ are week columns
    if months and num_cols > 1:
        col_offset = 1  # Skip label column
        for month in months:
            name = month.get("name", "")
            span = month.get("span", 1)
            if col_offset < num_cols:
                _set_cell_text(table.cell(0, col_offset), name, size=8, bold=True)
            col_offset += span

    # Row 1: Week numbers
    if weeks and num_rows > 1:
        for j, week in enumerate(weeks):
            col_idx = j + 1  # Offset for label column
            if col_idx < num_cols:
                _set_cell_text(table.cell(1, col_idx), week, size=7)

    # Rows 2+: Task rows
    if tasks and num_rows > 2:
        for i, task in enumerate(tasks):
            row_idx = i + 2  # Offset for header rows
            if row_idx >= num_rows:
                break

            task_name = task.get("name", "")
            cells = task.get("cells", [])
            task_type = task.get("type", "bar")  # "bar" (default) or "milestone"

            # Task name in first column
            _set_cell_text(table.cell(row_idx, 0), task_name, size=8,
                           bold=False, alignment=PP_ALIGN.LEFT)

            # Fill active cells based on task type
            for j in range(1, num_cols):
                cell_idx = j - 1  # 0-based cell index
                if cell_idx in cells:
                    if task_type == "milestone":
                        # Milestone rows (Kick-Off, Release Party): place a
                        # marker character instead of a colored bar fill
                        _set_cell_text(table.cell(row_idx, j), "\u25B6",
                                       size=7, color=GANTT_GREEN)
                    else:
                        # Sprint bar rows: fill cell with green background
                        _set_cell_fill(table.cell(row_idx, j), GANTT_GREEN)
                        _set_cell_text(table.cell(row_idx, j), "", size=7)

    if verbose:
        logger.info("Filled timeline/Gantt slide")


def fill_risks(slide, content: Dict, verbose: bool = False):
    """Fill progress/risks slide (slide 12) -- Bullet Points Lime layout.

    Two-column layout using existing placeholders:
    - PH idx=36: Left column header ("What already happened")
    - PH idx=42: Left column content
    - PH idx=46: Right column header ("Which risks do we see")
    - PH idx=45: Right column content
    """
    risks_data = content.get("risks", {})
    title = risks_data.get("title", "PROGRESS / RISKS")
    what_happened = risks_data.get("what_happened", [])
    risks = risks_data.get("risks", [])

    _set_placeholder_text(slide, 0, title.upper(), size=20, bold=True)
    title_ph = _find_placeholder(slide, 0)
    if title_ph and title_ph.has_text_frame:
        _enable_autofit(title_ph.text_frame)

    # Clear subtitle
    sub_ph = _find_placeholder(slide, 1)
    if sub_ph:
        _set_placeholder_text_theme(slide, 1, "")

    # Left column header (idx=36)
    _set_placeholder_text(slide, 36, "What already happened", size=14, bold=True,
                          color=OT_GREEN_ACCENT)

    # Left column content (idx=42)
    left_ph = _find_placeholder(slide, 42)
    if left_ph and left_ph.has_text_frame and what_happened:
        tf = left_ph.text_frame
        tf.clear()
        tf.word_wrap = True
        _enable_autofit(tf)
        for i, bullet in enumerate(what_happened):
            add_rich_paragraph(
                tf, bullet, font_size=11,
                use_theme_color=False, explicit_color=OT_ASH,
                green_color=OT_GREEN_HIGHLIGHT,
                spacing_before=3, spacing_after=3,
                is_first=(i == 0),
            )

    # Right column header (idx=46)
    _set_placeholder_text(slide, 46, "Which risks do we see", size=14, bold=True,
                          color=OT_GREEN_ACCENT)

    # Right column content (idx=45) -- sits on dark card (solidFill tx2)
    right_ph = _find_placeholder(slide, 45)
    if right_ph and right_ph.has_text_frame and risks:
        tf = right_ph.text_frame
        tf.clear()
        tf.word_wrap = True
        _enable_autofit(tf)
        for i, bullet in enumerate(risks):
            add_rich_paragraph(
                tf, bullet, font_size=11,
                use_theme_color=False, explicit_color=OT_LIME_TEXT,
                green_color=OT_GREEN_HIGHLIGHT,
                spacing_before=3, spacing_after=3,
                is_first=(i == 0),
            )

    # The template has this slide hidden (show="0").  Unhide it now that
    # we've filled content.
    if slide._element.get("show") == "0":
        slide._element.set("show", "1")

    if verbose:
        logger.info("Filled progress/risks slide")


def fill_participants(slide, content: Dict, verbose: bool = False):
    """Fill participants slide (slide 14) -- Bullet Points Lime layout.

    Elements:
    - idx=0: Title (e.g. "Our project has different participants")
    - PH idx=45: Text area for meeting type names (right side)
    - Template already contains colored rectangle shapes (Rechteck 40/41/42)
      with meeting type labels -- do NOT create or modify those shapes.
    - Template also contains a circular/hub diagram -- leave it untouched.
    """
    collab = content.get("collaboration", {})
    title = collab.get("participants_title", "PARTICIPANTS")
    meeting_types = collab.get("meeting_types", [])

    _set_placeholder_text(slide, 0, title.upper(), size=20, bold=True)
    title_ph = _find_placeholder(slide, 0)
    if title_ph and title_ph.has_text_frame:
        _enable_autofit(title_ph.text_frame)

    # Clear subtitle
    sub_ph = _find_placeholder(slide, 1)
    if sub_ph:
        _set_placeholder_text_theme(slide, 1, "")

    # Fill meeting types into placeholder idx=45 (right-side text area).
    # The template already provides colored rectangle shapes for the visual
    # layout -- we only fill the text placeholder, not create new shapes.
    ph45 = _find_placeholder(slide, 45)
    if ph45 and ph45.has_text_frame and meeting_types:
        tf = ph45.text_frame
        tf.clear()
        tf.word_wrap = True
        _enable_autofit(tf)

        for i, mt in enumerate(meeting_types):
            # meeting_types items can be plain strings or dicts with
            # "name" and optional "description" keys.
            if isinstance(mt, dict):
                mt_name = mt.get("name", "")
                mt_desc = mt.get("description", "")
            else:
                mt_name = str(mt)
                mt_desc = ""

            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(4)
            p.space_after = Pt(2) if mt_desc else Pt(4)
            r = p.add_run()
            r.text = mt_name
            r.font.size = Pt(12)
            r.font.name = "Akkurat LL"
            r.font.bold = True
            r.font.color.rgb = OT_ASH

            # Add description on a new line if provided
            if mt_desc:
                p_desc = tf.add_paragraph()
                p_desc.space_before = Pt(0)
                p_desc.space_after = Pt(4)
                r_desc = p_desc.add_run()
                r_desc.text = mt_desc
                r_desc.font.size = Pt(10)
                r_desc.font.name = "Akkurat LL"
                r_desc.font.bold = False
                r_desc.font.color.rgb = OT_ASH

    if verbose:
        logger.info("Filled participants slide")


def fill_meetings(slide, content: Dict, verbose: bool = False):
    """Fill meetings slide (slide 15) -- Bullet Points Lime layout.

    Two-column layout using existing placeholders/textboxes:
    - PH idx=36: Left meeting card header
    - PH idx=42: Left meeting card content
    - PH idx=46: Right meeting card header
    - Textbox id=7 at (~6.80, ~3.67): Right meeting card content
    """
    collab = content.get("collaboration", {})
    title = collab.get("meetings_title", "MEETINGS")
    meetings = collab.get("meetings", [])

    _set_placeholder_text(slide, 0, title.upper(), size=20, bold=True)
    title_ph = _find_placeholder(slide, 0)
    if title_ph and title_ph.has_text_frame:
        _enable_autofit(title_ph.text_frame)

    # Clear subtitle
    sub_ph = _find_placeholder(slide, 1)
    if sub_ph:
        _set_placeholder_text_theme(slide, 1, "")

    if not meetings:
        if verbose:
            logger.info("Filled meetings slide (no meetings data)")
        return

    def _format_meeting_text(meeting):
        name = meeting.get("name", "Meeting")
        day = meeting.get("day", "")
        time_val = meeting.get("time", "")
        frequency = meeting.get("frequency", "")
        participants = meeting.get("participants", [])
        lines = []
        if day:
            lines.append(f"DAY: {day}")
        if time_val:
            lines.append(f"TIME: {time_val}")
        if frequency:
            lines.append(f"FREQUENCY: {frequency}")
        if participants:
            lines.append(f"PARTICIPANTS: {', '.join(participants)}")
        return "\n".join(lines)

    # Select the 2 meetings to display on slide (template supports max 2).
    # When 3+ meetings, prioritize: Jour Fixe first, then Steering Committee.
    # All meetings (including overflow) are listed in speaker notes.
    if len(meetings) <= 2:
        display_meetings = list(meetings)
    else:
        # Priority order: Jour Fixe > Steering Committee > rest by original order
        priority_keywords = ["jour fixe", "steering"]
        prioritized = []
        remaining = []
        for m in meetings:
            name_lower = m.get("name", "").lower()
            matched = False
            for kw in priority_keywords:
                if kw in name_lower:
                    prioritized.append((priority_keywords.index(kw), m))
                    matched = True
                    break
            if not matched:
                remaining.append(m)
        prioritized.sort(key=lambda x: x[0])
        display_meetings = [m for _, m in prioritized] + remaining
        display_meetings = display_meetings[:2]

    def _fill_meeting_details(text_frame, meeting, text_color=OT_ASH):
        """Fill a text frame with formatted meeting detail lines."""
        tf = text_frame
        tf.clear()
        tf.word_wrap = True
        _enable_autofit(tf)
        detail_text = _format_meeting_text(meeting)
        for i, line in enumerate(detail_text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(2)
            p.space_after = Pt(2)
            r = p.add_run()
            r.text = line
            r.font.size = Pt(10)
            r.font.color.rgb = text_color

    # Meeting 1 -> Left column
    if len(display_meetings) >= 1:
        m = display_meetings[0]
        _set_placeholder_text(slide, 36, m.get("name", "Meeting 1"), size=14, bold=True,
                              color=OT_GREEN_ACCENT)

        left_ph = _find_placeholder(slide, 42)
        if left_ph and left_ph.has_text_frame:
            _fill_meeting_details(left_ph.text_frame, m, text_color=OT_ASH)

    # Meeting 2 -> Right column
    if len(display_meetings) >= 2:
        m = display_meetings[1]
        _set_placeholder_text(slide, 46, m.get("name", "Meeting 2"), size=14, bold=True,
                              color=OT_GREEN_ACCENT)

        # Right content is textbox id=7 at (6.80, 3.67)
        right_box = None
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                if shape.placeholder_format is not None:
                    continue
            except:
                pass
            l = (shape.left or 0) / 914400
            t = (shape.top or 0) / 914400
            w = (shape.width or 0) / 914400
            h = (shape.height or 0) / 914400
            if 6.0 < l < 7.5 and 3.0 < t < 4.5 and w > 4.0 and h > 2.0:
                right_box = shape
                break

        if right_box:
            _fill_meeting_details(right_box.text_frame, m, text_color=OT_LIME_TEXT)

    # Speaker notes: list ALL meetings so no data is lost.
    # When <=2 meetings, notes serve as a text reference.
    # When 3+ meetings, notes capture the overflow that can't fit on slide.
    if len(meetings) > 2:
        notes_text = "All project meetings:\n"
        for m in meetings:
            notes_text += f"\n{m.get('name', '')}"
            notes_text += f"\n  DAY: {m.get('day', '')} | TIME: {m.get('time', '')} | FREQUENCY: {m.get('frequency', '')}"
            parts = m.get('participants', [])
            if parts:
                notes_text += f"\n  PARTICIPANTS: {', '.join(parts)}"
        notes_text += f"\n\nNote: Only the 2 most important meetings are shown on the slide."
        add_speaker_notes(slide, notes_text)

    if verbose:
        logger.info("Filled meetings slide")


def fill_discussion(slide, content: Dict, verbose: bool = False):
    """Fill discussion slide (slide 17) -- Bullet Points Lime layout.

    Elements:
    - idx=0: Title
    - idx=1: Subtitle
    - Template already has a question text box (id=3) with discussion prompt
    """
    discussion = content.get("discussion", {})
    title = discussion.get("title", "OTHER TOPICS")
    subtitle = discussion.get("subtitle", "DISCUSSION")

    _set_placeholder_text(slide, 0, title.upper(), size=20, bold=True)
    title_ph = _find_placeholder(slide, 0)
    if title_ph and title_ph.has_text_frame:
        _enable_autofit(title_ph.text_frame)

    _set_placeholder_text(slide, 1, subtitle.upper(), size=14)
    sub_ph = _find_placeholder(slide, 1)
    if sub_ph and sub_ph.has_text_frame:
        _enable_autofit(sub_ph.text_frame)

    # Enable autofit on the existing discussion prompt text box (id=3 at ~(3.82, 3.27))
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            if shape.placeholder_format is not None:
                continue
        except Exception:
            pass
        l = (shape.left or 0) / 914400
        t = (shape.top or 0) / 914400
        w = (shape.width or 0) / 914400
        h = (shape.height or 0) / 914400
        if 3.0 < l < 5.0 and 2.5 < t < 4.0 and w > 5.0 and h > 1.0:
            _enable_autofit(shape.text_frame)
            # Reduce font size on the existing text to fit
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size > Pt(14):
                        run.font.size = Pt(14)
            break

    if verbose:
        logger.info("Filled discussion slide")


def fill_thankyou(slide, content: Dict, verbose: bool = False):
    """Fill thank-you slide (slide 19) -- Bullet Points Lime layout.

    Elements:
    - idx=0: Title ("THANK YOU!!")
    - Existing template text box for checkout question
    """
    checkout = content.get("checkout", {})
    title = checkout.get("title", "THANK YOU!!")
    question = checkout.get("question", "How do you feel after the kick-off?")

    _set_placeholder_text(slide, 0, title.upper(), size=24, bold=True)
    title_ph = _find_placeholder(slide, 0)
    if title_ph and title_ph.has_text_frame:
        _enable_autofit(title_ph.text_frame)

    # Clear subtitle
    sub_ph = _find_placeholder(slide, 1)
    if sub_ph:
        _set_placeholder_text_theme(slide, 1, "")

    # Find and update the existing question text box (id=3 at ~(2.04, 3.45))
    if question:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                if shape.placeholder_format is not None:
                    continue
            except Exception:
                pass
            l = (shape.left or 0) / 914400
            t = (shape.top or 0) / 914400
            w = (shape.width or 0) / 914400
            if 1.5 < l < 3.0 and 3.0 < t < 4.5 and w > 3.0:
                # Resize to give enough room for question text
                shape.height = Inches(0.80)
                tf = shape.text_frame
                tf.clear()
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                r = p.add_run()
                r.text = question
                r.font.size = Pt(14)
                r.font.color.rgb = OT_ASH
                _enable_autofit(tf)
                break

    if verbose:
        logger.info("Filled thank-you slide")


def _clear_layout_headline_prompts(prs):
    """Clear ALL custom prompt text from slide layout placeholders.

    The OT template's slide layouts have placeholders with hasCustomPrompt="1"
    and default text like 'Headline', 'Text', 'XX', '20XX', etc. When a slide
    doesn't override these placeholders, the prompt text bleeds through visually.
    Clearing ALL custom prompts at the layout level prevents this globally.

    For placeholders that also have visible fills (dark boxes), the fill is
    removed too so empty boxes don't render on slides.

    Slides that DO use these placeholders (e.g., Meetings, Risks, Timeline)
    already have them on the slide with real content, so clearing layout prompts
    is safe -- the slide-level placeholder always takes precedence.
    """
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }

    # Collect which placeholder indices are actually present on slides
    # (i.e., have been cloned from layout -> slide with real content).
    # We can safely hide layout-only placeholders.
    slide_ph_indices_by_layout = {}
    for slide in prs.slides:
        layout_name = slide.slide_layout.name if slide.slide_layout else None
        if layout_name not in slide_ph_indices_by_layout:
            slide_ph_indices_by_layout[layout_name] = set()
        for ph in slide.placeholders:
            slide_ph_indices_by_layout[layout_name].add(ph.placeholder_format.idx)

    cleared_count = 0
    for layout in prs.slide_layouts:
        layout_name = layout.name
        slide_indices = slide_ph_indices_by_layout.get(layout_name, set())

        for ph in layout.placeholders:
            if not ph.has_text_frame:
                continue

            ph_el = ph._element.find('.//p:ph', nsmap)
            if ph_el is None or not ph_el.get('hasCustomPrompt'):
                continue

            ph_idx = ph.placeholder_format.idx

            # Remove hasCustomPrompt so it won't show as prompt text
            del ph_el.attrib['hasCustomPrompt']

            # Clear all text runs in the placeholder
            for a_r in ph._element.findall('.//a:r', nsmap):
                a_t = a_r.find('a:t', nsmap)
                if a_t is not None:
                    a_t.text = ''

            # If this placeholder is NOT on any slide using this layout,
            # also hide it visually by removing its fill so empty boxes don't render
            if ph_idx not in slide_indices:
                sp_pr = ph._element.find(f'{{{NS_P}}}spPr')
                if sp_pr is not None:
                    # Remove any existing fill (DrawingML fill elements)
                    for child in list(sp_pr):
                        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                        if tag in ('solidFill', 'gradFill', 'pattFill', 'blipFill'):
                            sp_pr.remove(child)
                    # Add noFill so the shape is invisible
                    etree.SubElement(sp_pr, f"{{{NS_A}}}noFill")
                    # Also remove line/outline if present
                    for child in list(sp_pr):
                        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                        if tag == 'ln':
                            sp_pr.remove(child)

            cleared_count += 1

    if cleared_count:
        logger.debug(f"Cleared {cleared_count} layout custom prompts")


def update_copyright_footer(slide, copyright_year: str):
    """Update copyright footer textbox on a slide if it exists.

    Searches for ALL text boxes containing "ONE THOUSAND" and updates them.
    """
    copyright_text = f"\u00a9 {copyright_year} ONE THOUSAND"
    updated = False
    for shape in slide.shapes:
        if shape.has_text_frame:
            full_text = shape.text_frame.text.strip()
            if "ONE THOUSAND" in full_text or "\u00a9" in full_text:
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = copyright_text
                r.font.size = Pt(8)
                r.font.color.rgb = OT_GRAY
                logger.debug(f"Updated copyright footer on slide")
                updated = True
    return updated


def _update_layout_copyright(prs, copyright_year: str):
    """Update copyright footer in ALL slide layouts.

    Some slides (e.g. Thank You) don't have a slide-level copyright textbox
    and instead inherit from the layout. This ensures all layouts show the
    correct copyright year.
    """
    copyright_text = f"\u00a9 {copyright_year} ONE THOUSAND"
    updated = 0
    for layout in prs.slide_layouts:
        for shape in layout.shapes:
            if shape.has_text_frame:
                full_text = shape.text_frame.text.strip()
                if "ONE THOUSAND" in full_text or "\u00a9" in full_text:
                    tf = shape.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    r = p.add_run()
                    r.text = copyright_text
                    r.font.size = Pt(8)
                    r.font.color.rgb = OT_GRAY
                    updated += 1
    if updated:
        logger.debug(f"Updated {updated} layout copyright footers")


# ---------------------------------------------------------------------------
# Multi-UC Handling
# ---------------------------------------------------------------------------


def _duplicate_slide(prs, slide_index: int) -> int:
    """Duplicate a slide by copying its XML and relationships.

    Returns the index of the new slide (inserted right after the original).
    """
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from copy import deepcopy

    template_slide = prs.slides[slide_index]
    slide_layout = template_slide.slide_layout

    # Add a new slide with the same layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Copy all shapes from template to new slide
    # This is a simplified approach; for complex cases, XML-level copy is needed
    # For our purposes, the new slide gets filled by fill_pain_data separately

    return len(prs.slides) - 1


# ---------------------------------------------------------------------------
# Main Generation Pipeline
# ---------------------------------------------------------------------------


def generate_presentation(
    template_path: Path,
    variables_path: Path,
    content_path: Path,
    output_path: Path,
    verbose: bool = False,
):
    """Main entry point: load template, fill all 20 slides, save output.

    The template has 20 pre-existing slides. This function iterates through
    them and fills content into placeholders and text boxes based on the
    slide's position and layout.
    """
    variables = load_json(variables_path)
    content = load_json(content_path)

    logger.info(f"Loading template: {template_path}")
    prs = Presentation(str(template_path))

    # Clear "Headline" prompt text from slide layouts to prevent bleed-through
    _clear_layout_headline_prompts(prs)

    copyright_year = variables.get("copyright_year", "2019-2025")

    # Update copyright in slide layouts (some slides inherit from layout)
    _update_layout_copyright(prs, copyright_year)

    client_name = variables.get("client_name", "Client")
    project_title = variables.get("project_title", "Project")
    use_cases = variables.get("use_cases", [])
    images = variables.get("images", {})

    # Agenda sections from content
    agenda = content.get("agenda", {})
    sections = agenda.get("sections", [
        {"name": "Check In", "number": "01"},
        {"name": "Use-Case", "number": "02"},
        {"name": "Timeline", "number": "03"},
        {"name": "Collaboration", "number": "04"},
        {"name": "Other Topics", "number": "05"},
        {"name": "Check Out", "number": "06"},
    ])

    # Use-case data for multi-UC handling
    uc_content_list = content.get("use_cases", [])
    extra_uc_slides = max(0, len(uc_content_list) - 1)

    # Compute page numbers for agenda references
    page_numbers = _compute_page_numbers(len(prs.slides), extra_uc_slides)

    slides = list(prs.slides)
    num_slides = len(slides)
    logger.info(f"Template has {num_slides} slides")

    if num_slides < 20:
        logger.warning(
            f"Template has only {num_slides} slides, expected 20. "
            "Some slides may not be filled."
        )

    # -----------------------------------------------------------------------
    # Slide 0: Cover
    # -----------------------------------------------------------------------
    if num_slides > 0:
        fill_cover(slides[0], variables, content, verbose=verbose)

    # -----------------------------------------------------------------------
    # Slide 1: Full Agenda (ToC) -- no highlighting
    # -----------------------------------------------------------------------
    if num_slides > 1:
        fill_agenda(slides[1], sections, page_numbers,
                    active_section=0, slide_index=1, verbose=verbose)

    # -----------------------------------------------------------------------
    # Slide 2: Agenda Divider 01 -- Check In highlighted
    # -----------------------------------------------------------------------
    if num_slides > 2:
        fill_agenda(slides[2], sections, page_numbers,
                    active_section=1, slide_index=2, verbose=verbose)

    # -----------------------------------------------------------------------
    # Slide 3: Check-In
    # -----------------------------------------------------------------------
    if num_slides > 3:
        fill_checkin(slides[3], content, variables, images, verbose=verbose)

    # -----------------------------------------------------------------------
    # Slide 4: Agenda Divider 02 -- Use-case highlighted
    # -----------------------------------------------------------------------
    if num_slides > 4:
        fill_agenda(slides[4], sections, page_numbers,
                    active_section=2, slide_index=4, verbose=verbose)

    # -----------------------------------------------------------------------
    # Slide 5: Pain x Data (first use-case)
    # -----------------------------------------------------------------------
    if num_slides > 5:
        if uc_content_list:
            fill_pain_data(slides[5], uc_content_list[0], images, verbose=verbose)
        else:
            # Fallback: try to use flat use_case data
            uc_fallback = {
                "title": "PAIN x DATA",
                "pain_points": content.get("pain_points", []),
                "data_sources": content.get("data_sources", []),
                "solution": content.get("solution", []),
            }
            fill_pain_data(slides[5], uc_fallback, images, verbose=verbose)

    # -----------------------------------------------------------------------
    # Multi-UC: Duplicate slide 5 for additional use-cases
    # -----------------------------------------------------------------------
    if extra_uc_slides > 0 and num_slides > 5:
        for uc_idx in range(1, len(uc_content_list)):
            logger.info(f"Adding extra Pain x Data slide for UC {uc_idx + 1}")
            # Duplicate slide 5 layout by adding a new slide after current position
            layout_name = "Bullet Points Lime"
            layout = find_layout(prs, layout_name)
            if layout:
                new_slide = prs.slides.add_slide(layout)
                fill_pain_data(new_slide, uc_content_list[uc_idx], images,
                               verbose=verbose)
            else:
                logger.warning(f"Cannot duplicate Pain x Data: layout '{layout_name}' not found")

        # Recalculate page numbers after adding slides
        page_numbers = _compute_page_numbers(len(prs.slides), extra_uc_slides)

        # Re-fill all agenda slides with updated page numbers
        # (Only the original 20 slides are updated; extra UC slides are appended)
        slides = list(prs.slides)
        num_slides = len(slides)

        # Re-fill agenda slides
        agenda_slide_indices = [1, 2, 4]
        for si in agenda_slide_indices:
            if si < num_slides:
                active = DIVIDER_HIGHLIGHT_MAP.get(si, 0)
                _fill_agenda_slide(slides[si], sections, page_numbers, active)

    # -----------------------------------------------------------------------
    # Slide 6: Hackathon Validation
    # -----------------------------------------------------------------------
    if num_slides > 6:
        fill_hackathon(slides[6], content, images, verbose=verbose)

    # -----------------------------------------------------------------------
    # Slide 7: Step by Step
    # -----------------------------------------------------------------------
    if num_slides > 7:
        fill_step_by_step(slides[7], content, images, verbose=verbose)

    # -----------------------------------------------------------------------
    # Slide 8: Architecture
    # -----------------------------------------------------------------------
    if num_slides > 8:
        fill_architecture(slides[8], content, images, verbose=verbose)

    # -----------------------------------------------------------------------
    # Slide 9: Agenda Divider 03 -- Timeline highlighted
    # -----------------------------------------------------------------------
    if num_slides > 9:
        fill_agenda(slides[9], sections, page_numbers,
                    active_section=3, slide_index=9, verbose=verbose)

    # -----------------------------------------------------------------------
    # Slide 10: Sprint Goals
    # -----------------------------------------------------------------------
    if num_slides > 10:
        fill_sprint_goals(slides[10], content, verbose=verbose)

    # -----------------------------------------------------------------------
    # Slide 11: Timeline / Gantt
    # -----------------------------------------------------------------------
    if num_slides > 11:
        fill_timeline(slides[11], content, verbose=verbose)

    # -----------------------------------------------------------------------
    # Slide 12: Progress / Risks
    # -----------------------------------------------------------------------
    if num_slides > 12:
        fill_risks(slides[12], content, verbose=verbose)

    # -----------------------------------------------------------------------
    # Slide 13: Agenda Divider 04 -- Collaboration highlighted
    # -----------------------------------------------------------------------
    if num_slides > 13:
        fill_agenda(slides[13], sections, page_numbers,
                    active_section=4, slide_index=13, verbose=verbose)

    # -----------------------------------------------------------------------
    # Slide 14: Participants
    # -----------------------------------------------------------------------
    if num_slides > 14:
        fill_participants(slides[14], content, verbose=verbose)

    # -----------------------------------------------------------------------
    # Slide 15: Meetings
    # -----------------------------------------------------------------------
    if num_slides > 15:
        fill_meetings(slides[15], content, verbose=verbose)

    # -----------------------------------------------------------------------
    # Slide 16: Agenda Divider 05 -- Other Topics highlighted
    # -----------------------------------------------------------------------
    if num_slides > 16:
        fill_agenda(slides[16], sections, page_numbers,
                    active_section=5, slide_index=16, verbose=verbose)

    # -----------------------------------------------------------------------
    # Slide 17: Discussion
    # -----------------------------------------------------------------------
    if num_slides > 17:
        fill_discussion(slides[17], content, verbose=verbose)

    # -----------------------------------------------------------------------
    # Slide 18: Agenda Divider 06 -- Check Out highlighted
    # -----------------------------------------------------------------------
    if num_slides > 18:
        fill_agenda(slides[18], sections, page_numbers,
                    active_section=6, slide_index=18, verbose=verbose)

    # -----------------------------------------------------------------------
    # Slide 19: Thank You
    # -----------------------------------------------------------------------
    if num_slides > 19:
        fill_thankyou(slides[19], content, verbose=verbose)

    # -----------------------------------------------------------------------
    # Copyright footer update on all slides
    # -----------------------------------------------------------------------
    for slide in prs.slides:
        update_copyright_footer(slide, copyright_year)

    # -----------------------------------------------------------------------
    # Save output
    # -----------------------------------------------------------------------
    prs.save(str(output_path))
    total_slides = len(prs.slides)
    logger.info(f"Saved {total_slides} slides -> {output_path}")


# ---------------------------------------------------------------------------
# CLI Entry Point
# ---------------------------------------------------------------------------


def main():
    parser = argparse.ArgumentParser(
        description="Generate OT kick-off presentation from template + JSON inputs"
    )
    parser.add_argument(
        "--template", type=Path, required=True,
        help="Path to the kick-off template PPTX file"
    )
    parser.add_argument(
        "--variables", type=Path, required=True,
        help="Path to variables.json (client info, dates, images)"
    )
    parser.add_argument(
        "--content", type=Path, required=True,
        help="Path to content.json (slide content, agenda, bullets)"
    )
    parser.add_argument(
        "--output", type=Path, required=True,
        help="Path to write the output PPTX"
    )
    parser.add_argument(
        "--verbose", action="store_true",
        help="Enable verbose logging"
    )
    args = parser.parse_args()

    if args.verbose:
        logger.setLevel(logging.DEBUG)

    # Validate inputs
    if not args.template.exists():
        logger.error(f"Template not found: {args.template}")
        sys.exit(1)
    if not args.variables.exists():
        logger.error(f"Variables file not found: {args.variables}")
        sys.exit(1)
    if not args.content.exists():
        logger.error(f"Content file not found: {args.content}")
        sys.exit(1)

    # Ensure output directory exists
    output_dir = args.output.parent
    if not output_dir.exists():
        output_dir.mkdir(parents=True, exist_ok=True)
        logger.info(f"Created output directory: {output_dir}")

    try:
        generate_presentation(
            args.template, args.variables, args.content,
            args.output, args.verbose,
        )
    except Exception as e:
        logger.error(f"Generation failed: {e}")
        if args.verbose:
            import traceback
            traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
